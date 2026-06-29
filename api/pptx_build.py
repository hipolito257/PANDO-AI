"""
Vercel Python Serverless Function — PANDO PPTX Builder
POST /api/pptx_build
Body: { "template_url": str, "slide_plan": dict }
Returns: { "data": base64_pptx, "slide_count": int }
"""
import base64, io, json
from http.server import BaseHTTPRequestHandler

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData, XyChartData
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Constants ──────────────────────────────────────────────────────────────────
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DEFAULT_FONT = "Work Sans Light"

PALETTE = {
    "DKG": "004F46", "MDG": "437742", "OLV": "806E4B", "TEL": "4B5F62",
    "GRG": "D9DBD4", "NKB": "0A231F", "WHT": "FFFFFF", "LBL": "A5C8D1",
}

LAYOUT_MAP = {
    "takeaway": (1, 0),
    "divider":  (0, 2),
    "blank":    (2, 0),
}

PH = {"cat": 18, "title": 16, "takeaway": 17, "note": 14, "content": 26}


def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Builder ────────────────────────────────────────────────────────────────────
class PptxBuilder:
    def __init__(self, template_bytes: bytes):
        self.prs = Presentation(io.BytesIO(template_bytes))

    def build(self, slide_plan: dict) -> bytes:
        for sd in slide_plan.get("slides", []):
            self._add_slide(sd)
        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()

    def _add_slide(self, sd: dict):
        mi, li = LAYOUT_MAP.get(sd.get("layout", "takeaway"), (1, 0))
        layout = self.prs.slide_masters[mi].slide_layouts[li]
        slide = self.prs.slides.add_slide(layout)
        self._fill_phs(slide, sd)
        for el in sd.get("elements", []):
            self._element(slide, el)

    def _fill_phs(self, slide, sd: dict):
        mapping = {
            PH["cat"]: sd.get("category"),
            PH["title"]: sd.get("title"),
            PH["takeaway"]: sd.get("takeaway"),
            PH["note"]: sd.get("note"),
        }
        for ph in slide.placeholders:
            idx = ph.placeholder_format.idx
            val = mapping.get(idx)
            if val:
                ph.text = val
            elif idx == PH["content"]:
                ph._element.getparent().remove(ph._element)

    def _element(self, slide, el: dict):
        dispatch = {
            "panel_hdr":  self._panel_hdr,
            "textbox":    self._textbox,
            "shape":      self._shape,
            "hbar_float": self._hbar_float,
            "line":       self._line,
            "line_multi": self._line_multi,
            "donut":      self._donut,
            "scatter":    self._scatter,
            "quadrant":   self._quadrant,
        }
        fn = dispatch.get(el.get("type", ""))
        if fn:
            fn(slide, el)

    # ── Basic elements ─────────────────────────────────────────────────────────
    def _panel_hdr(self, slide, el: dict):
        x, y, w, h = el["x"], el["y"], el["w"], el.get("h", 0.27)
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(el.get("bg", PALETTE["DKG"]))
        sh.line.fill.background()
        p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = el.get("text", "")
        r.font.name = DEFAULT_FONT; r.font.size = Pt(el.get("size", 7.5))
        r.font.bold = True; r.font.color.rgb = _rgb(PALETTE["WHT"])

    def _textbox(self, slide, el: dict):
        box = slide.shapes.add_textbox(Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el.get("h", 0.25)))
        tf = box.text_frame; tf.word_wrap = el.get("wrap", True)
        align_map = {"c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT, "l": PP_ALIGN.LEFT}
        p = tf.paragraphs[0]; p.alignment = align_map.get(el.get("align", "l"), PP_ALIGN.LEFT)
        r = p.add_run(); r.text = el.get("text", "")
        r.font.name = DEFAULT_FONT; r.font.size = Pt(el.get("size", 7.5))
        r.font.bold = el.get("bold", False); r.font.italic = el.get("italic", False)
        r.font.color.rgb = _rgb(el.get("fg", PALETTE["NKB"]))

    def _shape(self, slide, el: dict):
        sh = slide.shapes.add_shape(1, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el.get("h", 0.27)))
        bg = el.get("bg")
        if bg: sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(bg)
        else: sh.fill.background()
        border = el.get("border")
        if border: sh.line.color.rgb = _rgb(border); sh.line.width = Pt(el.get("border_pt", 0.75))
        else: sh.line.fill.background()
        text = el.get("text", "")
        if text:
            p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = text
            r.font.name = DEFAULT_FONT; r.font.size = Pt(el.get("size", 8))
            r.font.bold = el.get("bold", False)
            r.font.color.rgb = _rgb(el.get("fg", PALETTE["WHT"]))

    # ── Chart helpers ──────────────────────────────────────────────────────────
    def _smooth_all(self, ch):
        for ser_el in ch._element.findall(".//{%s}ser" % C_NS):
            sm = ser_el.find("{%s}smooth" % C_NS)
            if sm is None: ser_el.append(parse_xml(f'<c:smooth xmlns:c="{C_NS}" val="1"/>'))
            else: sm.set("val", "1")

    def _fix_catax(self, ch, skip: int = 1):
        for catAx in ch._element.findall(".//{%s}catAx" % C_NS):
            noMulti  = catAx.find("{%s}noMultiLvlLbl" % C_NS)
            extLst_e = catAx.find("{%s}extLst" % C_NS)
            anchor = noMulti if noMulti is not None else (extLst_e if extLst_e is not None else None)
            for tag in ["tickLblSkip", "tickMarkSkip"]:
                el = catAx.find("{%s}%s" % (C_NS, tag))
                if el is None:
                    new_el = parse_xml(f'<c:{tag} xmlns:c="{C_NS}" val="{skip}"/>')
                    if anchor is not None: catAx.insert(list(catAx).index(anchor), new_el)
                    else: catAx.append(new_el)
                else:
                    el.set("val", str(skip))

    def _style_axes(self, ch, ymin=None, ymax=None, num_fmt="#,##0", skip=1, csize=5.5):
        va = ch.value_axis
        if ymin is not None: va.minimum_scale = ymin
        if ymax is not None: va.maximum_scale = ymax
        va.has_major_gridlines = True
        try:
            va.major_gridlines.format.line.color.rgb = _rgb("EBEBEB")
            va.major_gridlines.format.line.width = Pt(0.25)
        except Exception: pass
        va.tick_labels.font.size = Pt(6); va.tick_labels.font.color.rgb = _rgb("999999")
        va.tick_labels.number_format = num_fmt
        try: va.format.line.fill.background()
        except Exception: pass
        ca = ch.category_axis; ca.has_major_gridlines = False
        ca.tick_labels.font.size = Pt(csize); ca.tick_labels.font.color.rgb = _rgb("999999")
        try: ca.format.line.color.rgb = _rgb("DDDDD8"); ca.format.line.width = Pt(0.25)
        except Exception: pass
        self._fix_catax(ch, skip)

    def _set_line_style(self, series, color: str, width_pt: float = 1.5, dashed: bool = False):
        series.format.line.color.rgb = _rgb(color)
        series.format.line.width = Pt(width_pt)
        if dashed:
            ser_el = series._element
            spPr = ser_el.find("{%s}spPr" % C_NS)
            if spPr is not None:
                ln = spPr.find("{%s}ln" % A_NS)
                if ln is not None:
                    pd = etree.SubElement(ln, "{%s}prstDash" % A_NS)
                    pd.set("val", "dash")

    # ── Chart types ────────────────────────────────────────────────────────────
    def _hbar_float(self, slide, el: dict):
        series_defs = el.get("series", [])
        labels = [s["label"] for s in series_defs]
        lows   = [s["min"] for s in series_defs]
        highs  = [s["max"] - s["min"] for s in series_defs]
        cd = ChartData(); cd.categories = labels
        cd.add_series("spacer", lows); cd.add_series("range", highs)
        x, y, w, h = el["x"], el["y"], el["w"], el["h"]
        cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        sp = ch.series[0]; sp.format.fill.background(); sp.format.line.fill.background()
        colors = el.get("colors", [PALETTE["DKG"], PALETTE["MDG"], PALETTE["OLV"], PALETTE["TEL"], PALETTE["LBL"], PALETTE["GRG"]])
        self._color_bar_points(ch.series[1], colors, len(series_defs))
        ch.value_axis.has_major_gridlines = True
        try: ch.value_axis.major_gridlines.format.line.color.rgb = _rgb("EBEBEB"); ch.value_axis.major_gridlines.format.line.width = Pt(0.25)
        except Exception: pass
        ch.value_axis.tick_labels.font.size = Pt(6); ch.value_axis.tick_labels.number_format = "#,##0"
        ch.category_axis.has_major_gridlines = False; ch.category_axis.tick_labels.font.size = Pt(7)

    def _color_bar_points(self, series, colors: list, n: int):
        NS = f'xmlns:c="{C_NS}" xmlns:a="{A_NS}"'
        cat_el = series._element.find(qn("c:cat")) or series._element.find(qn("c:val"))
        idx = list(series._element).index(cat_el)
        for i in range(n):
            col = colors[i % len(colors)]
            series._element.insert(idx + i, parse_xml(
                f'<c:dPt {NS}><c:idx val="{i}"/><c:invertIfNegative val="0"/>'
                f'<c:spPr><a:solidFill><a:srgbClr val="{col}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></c:spPr></c:dPt>'
            ))

    def _line(self, slide, el: dict):
        cd = ChartData(); cd.categories = el["labels"]; cd.add_series("", el["values"])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        self._set_line_style(ch.series[0], el.get("color", PALETTE["MDG"]), el.get("width", 1.8))
        self._smooth_all(ch)
        self._style_axes(ch, ymin=el.get("ymin"), ymax=el.get("ymax"), num_fmt=el.get("num_fmt", "#,##0"), skip=el.get("skip", 6))

    def _line_multi(self, slide, el: dict):
        series_list = el.get("series", [])
        cd = ChartData(); cd.categories = el["labels"]
        for s in series_list:
            vals = s["values"]
            if len(vals) < len(el["labels"]):
                vals = vals + [vals[-1] if vals else 0] * (len(el["labels"]) - len(vals))
            cd.add_series(s["name"], vals)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]), cd)
        ch = cf.chart; ch.has_title = False
        for i, s in enumerate(series_list):
            self._set_line_style(ch.series[i], s.get("color", PALETTE["DKG"]), s.get("width", 1.4), dashed=s.get("dashed", False))
        self._smooth_all(ch)
        self._style_axes(ch, ymin=el.get("ymin", 0), ymax=el.get("ymax"), num_fmt=el.get("num_fmt", "#,##0"), skip=el.get("skip", 1), csize=5)
        ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False

    def _donut(self, slide, el: dict):
        slices = el.get("slices", [])
        cd = ChartData(); cd.categories = [s["label"] for s in slices]; cd.add_series("", [s["value"] for s in slices])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        ser = ch.series[0]; NS = f'xmlns:c="{C_NS}" xmlns:a="{A_NS}"'
        cat_el = ser._element.find(qn("c:cat")) or ser._element.find(qn("c:val"))
        idx = list(ser._element).index(cat_el)
        for i, s in enumerate(slices):
            col = s.get("color", PALETTE["GRG"])
            ser._element.insert(idx + i, parse_xml(
                f'<c:dPt {NS}><c:idx val="{i}"/>'
                f'<c:spPr><a:solidFill><a:srgbClr val="{col}"/></a:solidFill><a:ln><a:noFill/></a:ln></c:spPr></c:dPt>'
            ))
        hole = el.get("hole", 55)
        for dc in ch._element.findall(".//{%s}doughnutChart" % C_NS):
            hs = dc.find("{%s}holeSize" % C_NS)
            if hs is not None: hs.set("val", str(hole))
            else: dc.append(parse_xml(f'<c:holeSize xmlns:c="{C_NS}" val="{hole}"/>'))

    def _scatter(self, slide, el: dict):
        points = el.get("points", [])
        cd = XyChartData()
        for pt in points:
            s = cd.add_series(pt["label"]); s.add_data_point(pt["x"], pt["y"])
        cf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        for i, pt in enumerate(points):
            s = ch.series[i]; s.format.line.fill.background()
            s.marker.format.fill.solid(); s.marker.format.fill.fore_color.rgb = _rgb(pt.get("color", PALETTE["DKG"]))
            s.marker.format.line.fill.background(); s.marker.size = pt.get("size", 10)

    def _quadrant(self, slide, el: dict):
        x, y, w, h = el["x"], el["y"], el["w"], el["h"]
        mid_x = x + w / 2; mid_y = y + h / 2
        for conn, x1, y1, x2, y2 in [
            (slide.shapes.add_connector(1, Inches(mid_x), Inches(y), Inches(mid_x), Inches(y + h)), None, None, None, None),
            (slide.shapes.add_connector(1, Inches(x), Inches(mid_y), Inches(x + w), Inches(mid_y)), None, None, None, None),
        ]:
            conn.line.color.rgb = _rgb("CCCCCC"); conn.line.width = Pt(0.5)
        labels = el.get("axis_labels", {})
        def _lbl(text, lx, ly, lw=1.5, align="c"):
            box = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(lw), Inches(0.22))
            p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if align == "c" else PP_ALIGN.LEFT
            r = p.add_run(); r.text = text; r.font.name = DEFAULT_FONT; r.font.size = Pt(6.5); r.font.color.rgb = _rgb("888888")
        if labels.get("top"):    _lbl(labels["top"],    mid_x - 0.75, y - 0.25)
        if labels.get("bottom"): _lbl(labels["bottom"], mid_x - 0.75, y + h + 0.03)
        if labels.get("left"):   _lbl(labels["left"],   x - 1.6,      mid_y - 0.10)
        if labels.get("right"):  _lbl(labels["right"],  x + w + 0.05, mid_y - 0.10)
        for brand in el.get("brands", []):
            bx = x + brand["px"] * w; by = y + (1 - brand["py"]) * h
            col = brand.get("color", PALETTE["DKG"])
            dot = slide.shapes.add_shape(9, Inches(bx - 0.08), Inches(by - 0.08), Inches(0.16), Inches(0.16))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col); dot.line.fill.background()
            box = slide.shapes.add_textbox(Inches(bx - 0.5), Inches(by + 0.10), Inches(1.0), Inches(0.20))
            p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = brand.get("label", ""); r.font.name = DEFAULT_FONT; r.font.size = Pt(6.5); r.font.color.rgb = _rgb(col)


# ── Fetch template ─────────────────────────────────────────────────────────────
def fetch_template(url: str) -> bytes:
    if url.startswith("data:"):
        _, data = url.split(",", 1)
        return base64.b64decode(data)
    r = httpx.get(url, timeout=30, follow_redirects=True)
    r.raise_for_status()
    return r.content


# ── Vercel handler ─────────────────────────────────────────────────────────────
class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            body = json.loads(self.rfile.read(length))
            template_bytes = fetch_template(body["template_url"])
            builder = PptxBuilder(template_bytes)
            pptx_bytes = builder.build(body["slide_plan"])
            result = json.dumps({
                "data": base64.b64encode(pptx_bytes).decode(),
                "slide_count": len(body["slide_plan"].get("slides", [])),
            }).encode()
            self.send_response(200)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(result)))
            self.end_headers()
            self.wfile.write(result)
        except Exception as e:
            err = json.dumps({"error": str(e)}).encode()
            self.send_response(500)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(err)))
            self.end_headers()
            self.wfile.write(err)

    def log_message(self, *args):
        pass
