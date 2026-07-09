"""
PptxBuilder — builds PANDO presentations from a JSON slide plan.
Accepts a template PPTX (bytes) + slide_plan dict, returns PPTX bytes.
"""
import copy
import io
import re
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import ChartData, XyChartData
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from lxml import etree


# ── Constants ─────────────────────────────────────────────────────────────────
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DEFAULT_FONT = "Work Sans Light"

PALETTE = {
    "DKG": "004F46", "MDG": "437742", "OLV": "806E4B", "TEL": "4B5F62",
    "GRG": "D9DBD4", "NKB": "0A231F", "WHT": "FFFFFF", "LBL": "A5C8D1",
}

LAYOUT_MAP = {
    "cover":      (2, 0),
    "takeaway":   (1, 0),
    "divider":    (0, 2),
    "blank":      (2, 0),
    "back_cover": (2, 0),
}

PH = {"cat": 18, "title": 16, "takeaway": 17, "note": 14, "content": 26}

# Takeaway placeholder's real geometry in the template layout (idx=17): fixed
# top=1.078", height=0.551" (sized for up to ~2 lines at this font size). Used
# to compute exactly how far down its text actually reaches, instead of always
# assuming the box is full — see _content_start_y.
TAKEAWAY_TOP = 1.078
TAKEAWAY_SIZE = 14
TAKEAWAY_MAX_BOTTOM = 1.998  # the template's own fixed content-start y — never exceeded


# ── Safe value helpers ─────────────────────────────────────────────────────────

def _rgb(h) -> RGBColor:
    """Convert a hex string to RGBColor. Returns DKG green for any invalid input."""
    _fallback = RGBColor(0x00, 0x4F, 0x46)
    if not isinstance(h, str) or not h:
        return _fallback
    h = h.lstrip("#")
    if len(h) != 6:
        return _fallback
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, TypeError):
        return _fallback


def _str(val, default: str = "") -> str:
    """Return val as a string, falling back to default for None/False/True."""
    if val is None or val is False or val is True:
        return default
    return str(val)


def _num(val, default: float = 0.0) -> float:
    """Return val as float, falling back to default for non-numeric values."""
    try:
        return float(val)
    except (TypeError, ValueError):
        return default


def _tint(hex_color: str, amount: float) -> str:
    """Lighten a hex color toward white by `amount` (0 = unchanged, 1 = white).
    Used for subtle card background tints instead of hardcoding light colors,
    so comparison cards stay visually tied to the same palette as the charts."""
    h = (hex_color or "").lstrip("#")
    if len(h) != 6:
        return "F2F2F2"
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = round(r + (255 - r) * amount)
        g = round(g + (255 - g) * amount)
        b = round(b + (255 - b) * amount)
        return f"{r:02X}{g:02X}{b:02X}"
    except ValueError:
        return "F2F2F2"


CONTENT_X = 0.918      # left edge of title/takeaway/content placeholders in the template (839788 EMU)
CONTENT_RIGHT = 12.415  # right edge of the content placeholder (839788 + 10512714 EMU) — symmetric margin
CANVAS_BOTTOM = 7.15    # below this collides with the note/footer placeholder


def _geom(el: dict, defaults=(CONTENT_X, 2.0, 11.5, 4.0)):
    """Extract x, y, w, h from element dict with safe defaults and content-area clamps.
    y has a loose 1.3" safety floor only (never overlapping the title/category area) —
    the real, precise content-start position for a given slide (based on how much the
    takeaway actually wraps) is computed once in _add_slide and applied there by
    shifting every element's declared y, not by this floor.
    x minimum is 0.918" and right edge capped at 12.415" — matches the template's title/
    takeaway/content placeholder bounds exactly, so generated elements stay left- and
    right-aligned with the headers instead of drifting outside them."""
    x = max(_num(el.get("x"), defaults[0]), CONTENT_X)
    y = max(_num(el.get("y"), defaults[1]), 1.3)
    w = max(_num(el.get("w"), defaults[2]), 0.1)
    h = max(_num(el.get("h"), defaults[3]), 0.1)
    if x + w > CONTENT_RIGHT:
        w = max(CONTENT_RIGHT - x, 0.1)
    return x, y, w, h


# ── Builder ────────────────────────────────────────────────────────────────────

class PptxBuilder:
    def __init__(self, template_bytes: bytes, palette: dict | None = None, font: str | None = None):
        self.prs = Presentation(io.BytesIO(template_bytes))
        # Per-build palette/font: defaults to PANDO's own colors, but a profiled
        # template's theme (see template_profiler.py) overrides them so decks built
        # from a different uploaded template use THAT template's real colors/font
        # instead of Pando's, as long as the profile mapped at least the core keys.
        self.PALETTE = {**PALETTE, **{k: v for k, v in (palette or {}).items() if v}}
        self.DEFAULT_FONT = font or DEFAULT_FONT
        self.warnings: list[str] = []
        self._save_template_covers()
        self._clear_slides()

    def _pcolor(self, raw, default: str) -> str:
        """Snap a color value from the LLM-authored slide plan to the nearest
        PANDO brand hue. The build prompt tells the model to "never use colors
        outside the PANDO palette", but that's prose, not a schema constraint —
        models still sometimes emit an arbitrary off-brand hex (e.g. a gold
        badge color with no brand equivalent). This makes the constraint real:
        anything that isn't already an exact brand hex gets snapped to the
        closest one instead of rendering literally."""
        if not isinstance(raw, str) or not raw:
            return default
        key = raw.upper().lstrip("#")
        if key in self.PALETTE:
            return self.PALETTE[key]
        hexval = raw.lstrip("#")
        if len(hexval) != 6:
            return default
        try:
            r, g, b = int(hexval[0:2], 16), int(hexval[2:4], 16), int(hexval[4:6], 16)
        except ValueError:
            return default
        hexval = hexval.upper()
        brand_hexes = [self.PALETTE[k] for k in ("DKG", "MDG", "OLV", "TEL", "LBL", "GRG", "NKB", "WHT")]
        if hexval in brand_hexes:
            return hexval

        def _dist(h):
            br, bg, bb = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            return (r - br) ** 2 + (g - bg) ** 2 + (b - bb) ** 2

        return min(brand_hexes, key=_dist)

    def _save_template_covers(self):
        slides = list(self.prs.slides)
        self._cover_info = None
        self._back_cover_info = None
        if slides:
            self._cover_info = self._capture_slide(slides[0])
        if len(slides) >= 2:
            self._back_cover_info = self._capture_slide(slides[-1])

    def _capture_slide(self, slide) -> dict:
        rels = []
        for rId, rel in slide.part.rels.items():
            if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
                continue
            rels.append((rId, rel))
        return {"xml": copy.deepcopy(slide._element), "rels": rels}

    def _restore_slide(self, info: dict) -> Any:
        mi = min(0, len(self.prs.slide_masters) - 1)
        li = min(0, len(self.prs.slide_masters[mi].slide_layouts) - 1)
        layout = self.prs.slide_masters[mi].slide_layouts[li]
        slide = self.prs.slides.add_slide(layout)
        slide._element = info["xml"]
        for rId, rel in info["rels"]:
            try:
                slide.part._rels[rId] = rel
            except Exception:
                pass
        return slide

    def _clear_slides(self):
        sldIdLst = self.prs.slides._sldIdLst
        for sldId in list(sldIdLst):
            rId = sldId.get(qn("r:id"))
            try:
                self.prs.slides.part.drop_rel(rId)
            except Exception:
                pass
            sldIdLst.remove(sldId)

    # ── Public ─────────────────────────────────────────────────────────────────
    def build(self, slide_plan: dict) -> bytes:
        slides_data = slide_plan.get("slides", [])
        cover_data      = next((s for s in slides_data if s.get("layout") == "cover"),      {})
        back_cover_data = next((s for s in slides_data if s.get("layout") == "back_cover"), {})
        content_slides  = [s for s in slides_data if s.get("layout") not in ("cover", "back_cover")]

        if self._cover_info:
            cover_slide = self._restore_slide(self._cover_info)
            if cover_data:
                self._update_cover_text(cover_slide, cover_data)
        else:
            mi = min(2, len(self.prs.slide_masters) - 1)
            li = min(0, len(self.prs.slide_masters[mi].slide_layouts) - 1)
            slide = self.prs.slides.add_slide(self.prs.slide_masters[mi].slide_layouts[li])
            self._draw_cover(slide, cover_data)

        for sd in content_slides:
            try:
                self._add_slide(sd)
            except Exception:
                pass  # skip broken slides rather than crashing the build

        if self._back_cover_info:
            self._restore_slide(self._back_cover_info)
        else:
            mi = min(2, len(self.prs.slide_masters) - 1)
            li = min(0, len(self.prs.slide_masters[mi].slide_layouts) - 1)
            slide = self.prs.slides.add_slide(self.prs.slide_masters[mi].slide_layouts[li])
            self._draw_back_cover(slide, back_cover_data)

        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()

    def _update_cover_text(self, slide, cover_data: dict):
        title    = _str(cover_data.get("title")).strip()
        subtitle = _str(cover_data.get("subtitle")).strip()
        if not title and not subtitle:
            return

        # Rank by known run font size when the placeholder already has text (e.g.
        # a restored slide with sample copy); most templates' cover placeholders
        # are empty until now, so fall back to shape area (bigger box == title).
        candidates = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            max_sz = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = run.font.size or 0
                    if sz > max_sz:
                        max_sz = sz
            area = (shape.width or 0) * (shape.height or 0)
            candidates.append((max_sz, area, shape))
        candidates.sort(key=lambda c: (c[0], c[1]), reverse=True)
        shapes_by_size = [(sz, shape) for sz, _area, shape in candidates]

        def _replace_text(shape, new_text: str):
            tf = shape.text_frame
            first_run_el = None
            for para in tf.paragraphs:
                for run in para.runs:
                    first_run_el = copy.deepcopy(run._r)
                    break
                if first_run_el is not None:
                    break
            tf.clear()
            p = tf.paragraphs[0]
            if first_run_el is not None:
                t_el = first_run_el.find(qn("a:t"))
                if t_el is not None:
                    t_el.text = new_text
                p._p.append(first_run_el)
                # The cloned run keeps the template's own rPr (bold/italic/size),
                # which is exactly what we want — but the template's cover runs
                # often have no explicit <a:latin> override at all, so they'd
                # silently inherit the theme's minor font (Aptos) instead of
                # PANDO's brand font. Force it explicitly.
                p.runs[-1].font.name = self.DEFAULT_FONT
            else:
                r = p.add_run()
                r.text = new_text
                r.font.name = self.DEFAULT_FONT
                r.font.color.rgb = _rgb(self.PALETTE["WHT"])

        if title and shapes_by_size:
            _replace_text(shapes_by_size[0][1], title)
        if subtitle and len(shapes_by_size) > 1:
            title_shape = shapes_by_size[0][1]
            for _, shape in shapes_by_size[1:]:
                if shape is not title_shape:
                    _replace_text(shape, subtitle)
                    break

    # ── Precise content-start calculation ───────────────────────────────────────
    def _content_start_y(self, sd: dict) -> float:
        """Where the first content element should start, computed from the
        takeaway's actual wrapped height at TAKEAWAY_SIZE instead of always
        assuming its box is fully used. The takeaway placeholder's box is fixed
        (top=1.078", height=0.551", sized for ~2 lines) and independent of the
        title above it (the title has its own auto-shrink-to-fit, so it can never
        push into the takeaway's position) — only the takeaway's own text length
        affects how much of that box is actually filled.
        """
        takeaway = _str(sd.get("takeaway"))
        if not takeaway:
            # No takeaway at all — content can start right after the title's box.
            return TAKEAWAY_TOP
        col_w = CONTENT_RIGHT - CONTENT_X
        lines = min(self._wrap_lines(takeaway, col_w, TAKEAWAY_SIZE), 2)  # box fits ~2 lines; more just shrinks font, not height
        line_h = (TAKEAWAY_SIZE * 1.25) / 72
        bottom = TAKEAWAY_TOP + lines * line_h + 0.14  # + top/bottom internal padding
        return min(bottom + 0.20, TAKEAWAY_MAX_BOTTOM)  # +breathing room, never past the template's own max

    # ── Slide assembly ─────────────────────────────────────────────────────────
    def _add_slide(self, sd: dict):
        layout_key = _str(sd.get("layout"), "takeaway")
        mi, li = LAYOUT_MAP.get(layout_key, (1, 0))
        mi = min(mi, len(self.prs.slide_masters) - 1)
        li = min(li, len(self.prs.slide_masters[mi].slide_layouts) - 1)
        layout = self.prs.slide_masters[mi].slide_layouts[li]
        slide = self.prs.slides.add_slide(layout)

        self._fill_phs(slide, sd)

        elements = [el for el in sd.get("elements", []) if isinstance(el, dict)]
        # Shift all elements so the topmost one starts exactly where the takeaway's
        # actual text ends (see _content_start_y) — not a blind fixed 2.0" assumption.
        # A short one-line takeaway used to leave nearly half an inch of dead space
        # before the first chart; a long one relied on auto-shrink and the old fixed
        # value to avoid overlap. This corrects both directions in one calculation.
        if elements and layout_key != "divider":
            content_start = self._content_start_y(sd)
            min_y = min(_num(el.get("y"), content_start) for el in elements)
            shift = min_y - content_start
            if abs(shift) > 0.01:
                for el in elements:
                    el["y"] = _num(el.get("y"), min_y) - shift

        # Tables have a physical minimum row height (font size + padding) that
        # PowerPoint enforces regardless of the "h" the plan declared — a table
        # with more rows than the declared h assumed will silently grow taller
        # than the plan says, overlapping whatever sits below it. Recompute the
        # real height from the row count up front so QA and layout both see it,
        # and push down anything else in the table's column that was placed
        # assuming the plan's (too-small) height — otherwise a heading or card
        # meant to sit just below the table renders on top of its last row.
        for el in elements:
            if _str(el.get("type")) == "table":
                old_h = _num(el.get("h"))
                new_h = self._table_height(el)
                el["h"] = new_h
                delta = new_h - old_h
                if delta > 0.01:
                    tx, ty, tw = _num(el.get("x")), _num(el.get("y")), _num(el.get("w"))
                    old_bottom = ty + old_h
                    for other in elements:
                        if other is el:
                            continue
                        ox, oy, ow = _num(other.get("x")), _num(other.get("y")), _num(other.get("w"))
                        overlaps_x = ox < tx + tw and ox + ow > tx
                        if overlaps_x and oy >= old_bottom - 0.05:
                            other["y"] = oy + delta

        # The plan's y/h numbers are LLM-estimated and occasionally run past the
        # note/footer placeholder near the bottom of the slide (e.g. a chart given
        # h=2.05 at y=5.5 when the footer sits at y=6.84) — clamp so every visual
        # element's own box stays within the printable canvas rather than silently
        # rendering through the footer.
        for el in elements:
            if _str(el.get("type")) not in self._VISUAL_TYPES:
                continue
            y = _num(el.get("y"))
            h = _num(el.get("h"))
            if h and y + h > CANVAS_BOTTOM + 0.02:
                el["h"] = max(0.3, CANVAS_BOTTOM - y)

        self._qa_check_slide(sd, elements)

        for el in elements:
            self._element(slide, el)

    # ── Geometry-based QA (no rendering required) ──────────────────────────────
    # Catches the two most common AI-generated-deck defects — elements placed
    # off the printable canvas, and two visual elements overlapping — by pure
    # bounding-box math over the plan's own x/y/w/h. This can't catch text
    # overflow inside a shape (that needs an actual renderer), but it runs on
    # every build with no external dependency and reliably flags layout bugs
    # before they reach the user.
    _VISUAL_TYPES = {
        "bar", "line", "line_multi", "hbar_float", "donut", "scatter", "quadrant",
        "table", "stat_row", "icon_row", "comparison_cards", "timeline", "waterfall",
        "alt_timeline", "org_chart", "process_flow", "pill_row",
    }

    def _qa_check_slide(self, sd: dict, elements: list[dict]):
        title = _str(sd.get("title"), "(untitled slide)")
        top_bound = self._content_start_y(sd) - 0.05 if _str(sd.get("layout"), "takeaway") != "divider" else 1.95
        canvas = (CONTENT_X - 0.05, top_bound, CONTENT_RIGHT + 0.05, CANVAS_BOTTOM)

        boxes = []
        for el in elements:
            t = _str(el.get("type"))
            x, y, w, h = _geom(el)
            if x < canvas[0] - 0.02 or y < canvas[1] - 0.02 or (x + w) > canvas[2] + 0.02 or (y + h) > canvas[3] + 0.02:
                self.warnings.append(
                    f'Slide "{title}": {t} element at ({x:.2f},{y:.2f},{w:.2f}x{h:.2f}) extends outside the printable canvas.'
                )
            if t in self._VISUAL_TYPES:
                boxes.append((t, x, y, w, h))

        for i in range(len(boxes)):
            t1, x1, y1, w1, h1 = boxes[i]
            for j in range(i + 1, len(boxes)):
                t2, x2, y2, w2, h2 = boxes[j]
                ox = min(x1 + w1, x2 + w2) - max(x1, x2)
                oy = min(y1 + h1, y2 + h2) - max(y1, y2)
                if ox <= 0 or oy <= 0:
                    continue
                overlap_area = ox * oy
                smaller_area = min(w1 * h1, w2 * h2)
                if smaller_area > 0 and overlap_area / smaller_area > 0.15:
                    self.warnings.append(
                        f'Slide "{title}": {t1} and {t2} elements overlap by '
                        f'{overlap_area / smaller_area * 100:.0f}% of the smaller element\'s area.'
                    )

    def _fill_phs(self, slide, sd: dict):
        layout_key = _str(sd.get("layout"), "takeaway")
        title_val    = _str(sd.get("title"))    or None
        takeaway_val = _str(sd.get("takeaway")) or None
        cat_val      = _str(sd.get("category")) or None
        note_val     = _str(sd.get("note"))     or None

        phs = list(slide.placeholders)

        if layout_key == "divider":
            # Divider slides live on a different master with different placeholder
            # indices — fill by vertical position rather than by index number.
            text_phs = sorted(
                [p for p in phs if p.has_text_frame],
                key=lambda p: p.top if p.top is not None else 0,
            )
            filled = set()
            if text_phs and title_val:
                self._set_markdown_text(text_phs[0], title_val)
                filled.add(id(text_phs[0]))
            if len(text_phs) > 1 and takeaway_val:
                self._set_markdown_text(text_phs[1], takeaway_val)
                filled.add(id(text_phs[1]))
            for ph in phs:
                if id(ph) not in filled:
                    sp = ph._element
                    parent = sp.getparent()
                    if parent is not None:
                        parent.remove(sp)
            return

        # Regular content slides: fill by known placeholder index
        mapping = {
            PH["cat"]:      cat_val,
            PH["title"]:    title_val,
            PH["takeaway"]: takeaway_val,
            PH["note"]:     note_val,
        }
        # Takeaway is explicitly bumped to TAKEAWAY_SIZE (14pt) — the template's
        # own placeholder default is 12pt, sized for our house style.
        sizes = {PH["takeaway"]: TAKEAWAY_SIZE}
        to_remove = []
        for ph in phs:
            idx = ph.placeholder_format.idx
            val = mapping.get(idx)
            if val:
                self._set_markdown_text(ph, val, font_size=sizes.get(idx))
            else:
                to_remove.append(ph)
        for ph in to_remove:
            sp = ph._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)

    # ── Cover fallback ─────────────────────────────────────────────────────────
    def _draw_cover(self, slide, sd: dict):
        W, H = 13.33, 7.5

        def _txt(text, x, y, w, h, size, bold=False, italic=False, fg=self.PALETTE["NKB"], align="l", wrap=True):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = box.text_frame; tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            r = p.add_run(); r.text = _str(text)
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
            r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = _rgb(fg)

        def _rect(x, y, w, h, color):
            sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(color)
            sh.line.fill.background()

        _rect(0, 0, 0.40, H, self.PALETTE["DKG"])
        _rect(0.40, H * 0.6, W - 0.40, 0.04, self.PALETTE["OLV"])
        _txt(sd.get("title", sd.get("company", "")), 0.75, 2.20, W - 1.2, 1.40, 44, bold=True, fg=self.PALETTE["NKB"])
        _txt(sd.get("subtitle", "Investment Overview"),  0.75, 3.80, W - 1.2, 0.60, 20, fg=self.PALETTE["TEL"])
        _rect(0.75, 3.70, W - 1.5, 0.025, self.PALETTE["DKG"])
        _txt("Private & Confidential", 0.75, H - 0.65, 6, 0.35, 8, italic=True, fg="999999")
        _txt("STRICTLY CONFIDENTIAL", 0.02, 2.5, 0.28, 3.5, 6.5, fg=self.PALETTE["WHT"], align="c", wrap=True)

    def _draw_back_cover(self, slide, sd: dict):
        W, H = 13.33, 7.5

        def _rect(x, y, w, h, color):
            sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(color)
            sh.line.fill.background()

        def _txt(text, x, y, w, h, size, bold=False, italic=False, fg=self.PALETTE["WHT"], align="c"):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = box.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            r = p.add_run(); r.text = _str(text)
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
            r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = _rgb(fg)

        _rect(0, 0, W, H, self.PALETTE["DKG"])
        _rect(0, H * 0.75, W, 0.05, self.PALETTE["OLV"])
        _txt(sd.get("title", "Preguntas"), 1, H / 2 - 0.8, W - 2, 1.4, 48, bold=True)
        subtitle = _str(sd.get("subtitle"))
        if subtitle:
            _txt(subtitle, 1, H / 2 + 0.7, W - 2, 0.5, 16, fg="A5C8D1")
        _txt("pando.vc  |  Private & Confidential", 0, H - 0.55, W, 0.35, 9, italic=True)

    def _set_markdown_text(self, ph, text: str, font_size: float | None = None):
        tf = ph.text_frame
        tf.clear()
        lines = _str(text).split("\n")
        for line_idx, line in enumerate(lines):
            p = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
            for part in re.split(r'(\*\*.*?\*\*)', line):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**"):
                    r.text = part[2:-2]
                    r.font.bold = True
                else:
                    r.text = part
                # Force the brand font explicitly rather than letting these
                # placeholders (title/takeaway/category/note) fall back to
                # whatever the template's own placeholder style happens to be —
                # that's how the footnote/source line ended up in plain "Work
                # Sans" instead of "Work Sans Light" while every other element
                # in this file explicitly sets DEFAULT_FONT.
                r.font.name = self.DEFAULT_FONT
                if font_size is not None:
                    r.font.size = Pt(font_size)
        # A long title/takeaway that wraps to a second line can overflow its
        # placeholder's fixed box and collide with the text below it (e.g. the
        # takeaway paragraph sitting right under the title). Let PowerPoint
        # shrink the text to fit rather than overflow — it only kicks in when
        # the text actually doesn't fit, so single-line titles are unaffected.
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

    # ── Element dispatcher ─────────────────────────────────────────────────────
    def _element(self, slide, el: dict):
        t = _str(el.get("type"))

        # Skip text-only elements with no actual content
        if t in ("textbox", "panel_hdr") and not _str(el.get("text")).strip():
            return
        if t == "shape" and not _str(el.get("text")).strip() and not el.get("bg") and not el.get("border"):
            return
        # Skip chart elements with no data
        if t in ("line", "line_multi", "bar") and not el.get("labels") and not el.get("series"):
            return
        if t == "hbar_float" and not el.get("series"):
            return
        if t == "donut" and not el.get("slices"):
            return
        if t == "scatter" and not el.get("points"):
            return
        if t == "table" and not el.get("headers"):
            return
        if t == "stat_row" and not el.get("items"):
            return
        if t == "icon_row" and not el.get("items"):
            return
        if t == "comparison_cards" and not el.get("cards"):
            return
        if t == "timeline" and not el.get("steps"):
            return
        if t == "waterfall" and not el.get("labels"):
            return
        if t == "alt_timeline" and not el.get("entries"):
            return
        if t == "org_chart" and not el.get("levels"):
            return
        if t == "process_flow" and not el.get("steps"):
            return
        if t == "pill_row" and not el.get("items"):
            return
        if t == "band_scatter" and not el.get("points"):
            return

        dispatch = {
            "panel_hdr":  self._panel_hdr,
            "textbox":    self._textbox,
            "shape":      self._shape,
            "hbar_float": self._hbar_float,
            "line":       self._line,
            "line_multi": self._line_multi,
            "bar":        self._bar,
            "donut":      self._donut,
            "scatter":    self._scatter,
            "quadrant":   self._quadrant,
            "table":      self._table,
            "stat_row":         self._stat_row,
            "icon_row":         self._icon_row,
            "comparison_cards": self._comparison_cards,
            "timeline":         self._timeline,
            "waterfall":        self._waterfall,
            "alt_timeline":     self._alt_timeline,
            "org_chart":        self._org_chart,
            "process_flow":     self._process_flow,
            "pill_row":         self._pill_row,
            "band_scatter":     self._band_scatter,
        }
        fn = dispatch.get(t)
        if not fn:
            return
        try:
            fn(slide, el)
        except Exception:
            pass  # skip broken elements; never crash the whole slide

    # ── Basic shapes / text ────────────────────────────────────────────────────
    def _panel_hdr(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 0.27))
        h = max(_num(el.get("h"), 0.27), 0.1)
        bg = el.get("bg", self.PALETTE["DKG"])
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(bg)
        sh.line.fill.background()
        tf = sh.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = _str(el.get("text"))
        r.font.name = self.DEFAULT_FONT; r.font.size = Pt(_num(el.get("size"), 7.5))
        r.font.bold = True; r.font.color.rgb = _rgb(self.PALETTE["WHT"])

    def _textbox(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 5.0, 0.25))
        h = max(_num(el.get("h"), 0.25), 0.1)
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = el.get("wrap", True)
        align_map = {"c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT, "l": PP_ALIGN.LEFT}
        font_size  = Pt(_num(el.get("size"), 7.5))
        font_bold  = bool(el.get("bold", False))
        font_ital  = bool(el.get("italic", False))
        font_color = _rgb(el.get("fg", self.PALETTE["NKB"]))
        align      = align_map.get(_str(el.get("align"), "l"), PP_ALIGN.LEFT)
        lines = _str(el.get("text")).split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            for part in re.split(r'(\*\*.*?\*\*)', line):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    r.text = part[2:-2]; r.font.bold = True
                else:
                    r.text = part; r.font.bold = font_bold
                r.font.name = self.DEFAULT_FONT; r.font.size = font_size
                r.font.italic = font_ital
                r.font.color.rgb = font_color
        # Free-form textboxes are common for KPI-style callouts (e.g. a big
        # "~15,000" stat in a narrow column) where the plan's font size doesn't
        # always fit the declared width — shrink rather than wrap/overflow.
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass

    def _shape(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 5.0, 0.27))
        h = max(_num(el.get("h"), 0.27), 0.1)
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        bg = el.get("bg")
        if bg and isinstance(bg, str):
            sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(bg)
        else:
            sh.fill.background()
        border = el.get("border")
        if border and isinstance(border, str):
            sh.line.color.rgb = _rgb(border)
            sh.line.width = Pt(_num(el.get("border_pt"), 0.75))
        else:
            sh.line.fill.background()
        text = _str(el.get("text"))
        if text:
            tf = sh.text_frame; tf.word_wrap = True
            font_size  = Pt(_num(el.get("size"), 8))
            font_bold  = bool(el.get("bold", False))
            font_color = _rgb(el.get("fg", self.PALETTE["WHT"]))
            lines = text.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = line
                r.font.name = self.DEFAULT_FONT; r.font.size = font_size
                r.font.bold = font_bold; r.font.color.rgb = font_color

    # ── Charts: shared helpers ─────────────────────────────────────────────────
    def _smooth_all(self, ch):
        for ser_el in ch._element.findall(".//{%s}ser" % C_NS):
            sm = ser_el.find("{%s}smooth" % C_NS)
            if sm is None:
                ser_el.append(parse_xml(f'<c:smooth xmlns:c="{C_NS}" val="1"/>'))
            else:
                sm.set("val", "1")

    def _fix_catax(self, ch, skip: int = 1):
        for catAx in ch._element.findall(".//{%s}catAx" % C_NS):
            noMulti  = catAx.find("{%s}noMultiLvlLbl" % C_NS)
            extLst_e = catAx.find("{%s}extLst" % C_NS)
            anchor = noMulti if noMulti is not None else (extLst_e if extLst_e is not None else None)
            for tag in ["tickLblSkip", "tickMarkSkip"]:
                el = catAx.find("{%s}%s" % (C_NS, tag))
                if el is None:
                    new_el = parse_xml(f'<c:{tag} xmlns:c="{C_NS}" val="{skip}"/>')
                    if anchor is not None:
                        catAx.insert(list(catAx).index(anchor), new_el)
                    else:
                        catAx.append(new_el)
                else:
                    el.set("val", str(skip))

    def _set_chart_default_font(self, ch):
        """python-pptx always writes a chartSpace-level c:txPr fallback (18pt,
        no typeface) that any unstyled chart text inherits from the theme's
        minor font (Aptos) instead of our Work Sans Light — patch it in so
        nothing inside a chart can silently fall back to the theme font."""
        try:
            defRPr = ch._chartSpace.find(qn("c:txPr")).find(qn("a:p")).find(qn("a:pPr")).find(qn("a:defRPr"))
            if defRPr is not None and defRPr.find(qn("a:latin")) is None:
                latin = etree.SubElement(defRPr, qn("a:latin"))
                latin.set("typeface", self.DEFAULT_FONT)
        except Exception: pass

    def _style_axes(self, ch, ymin=None, ymax=None, num_fmt="#,##0",
                    skip=1, csize=5.5, grid_color="EBEBEB"):
        try:
            va = ch.value_axis
            if ymin is not None: va.minimum_scale = _num(ymin)
            if ymax is not None: va.maximum_scale = _num(ymax)
            va.has_major_gridlines = True
            try:
                va.major_gridlines.format.line.color.rgb = _rgb(grid_color)
                va.major_gridlines.format.line.width = Pt(0.25)
            except Exception: pass
            va.tick_labels.font.size = Pt(6)
            va.tick_labels.font.name = self.DEFAULT_FONT
            va.tick_labels.font.color.rgb = _rgb("999999")
            va.tick_labels.number_format = num_fmt
            try: va.format.line.fill.background()
            except Exception: pass
        except Exception: pass
        try:
            ca = ch.category_axis
            ca.has_major_gridlines = False
            ca.tick_labels.font.size = Pt(csize)
            ca.tick_labels.font.name = self.DEFAULT_FONT
            ca.tick_labels.font.color.rgb = _rgb("999999")
            try:
                ca.format.line.color.rgb = _rgb("DDDDD8")
                ca.format.line.width = Pt(0.25)
            except Exception: pass
            self._fix_catax(ch, skip)
        except Exception: pass

    def _plot_rect(self, x: float, y: float, w: float, h: float,
                  left: float = 0.42, right: float = 0.08, top: float = 0.05, bottom: float = 0.32):
        """Approximate the plotted data area inside a chart's outer box (reserving
        margin for axis tick labels) — used to position manually-drawn overlay
        annotations (pps-delta call-outs, category highlights) that native
        chart XML can't express, without needing pixel-exact chart internals."""
        return x + left, y + top, max(w - left - right, 0.1), max(h - top - bottom, 0.1)

    def _box_data_labels(self, ser, border_color: str = "444444", fill_color: str = "FFFFFF"):
        """Give a series' data labels a thin-bordered white box behind the number —
        the signature PANDO chart callout (e.g. '1,727', '80%', '37%' boxed above
        each bar/point) instead of PowerPoint's plain floating label text."""
        try:
            dLbls = ser.data_labels._element
            if dLbls.find(qn("c:spPr")) is not None:
                return
            spPr = parse_xml(
                f'<c:spPr xmlns:c="{C_NS}" xmlns:a="{A_NS}">'
                f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'
                f'<a:ln w="{int(Pt(0.75))}"><a:solidFill><a:srgbClr val="{border_color}"/></a:solidFill></a:ln>'
                f'</c:spPr>'
            )
            numFmt = dLbls.find(qn("c:numFmt"))
            if numFmt is not None:
                numFmt.addnext(spPr)
            else:
                dLbls.insert(0, spPr)
        except Exception: pass

    def _style_legend(self, ch):
        """Small bottom legend — PowerPoint's default legend text is huge (~18pt)
        next to our 6-7pt axis labels, so always shrink it to match the deck."""
        try:
            ch.has_legend = True
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
            ch.legend.font.size = Pt(7.5)
            ch.legend.font.name = self.DEFAULT_FONT
            ch.legend.font.color.rgb = _rgb("444444")
        except Exception: pass

    def _set_line_style(self, series, color, width_pt: float = 1.5, dashed: bool = False):
        try:
            series.format.line.color.rgb = _rgb(color)
            series.format.line.width = Pt(_num(width_pt, 1.5))
            if dashed:
                ser_el = series._element
                spPr = ser_el.find("{%s}spPr" % C_NS)
                if spPr is not None:
                    ln = spPr.find("{%s}ln" % A_NS)
                    if ln is not None:
                        pd = etree.SubElement(ln, "{%s}prstDash" % A_NS)
                        pd.set("val", "dash")
        except Exception: pass

    def _color_bar_points(self, series, colors: list, n: int):
        NS = (f'xmlns:c="{C_NS}" xmlns:a="{A_NS}"')
        cat_el = series._element.find(qn("c:cat"))
        if cat_el is None:
            cat_el = series._element.find(qn("c:val"))
        if cat_el is None:
            return  # can't determine insertion point — skip coloring
        idx = list(series._element).index(cat_el)
        for i in range(n):
            col = _str(colors[i % len(colors)], self.PALETTE["DKG"]) if colors else self.PALETTE["DKG"]
            series._element.insert(idx + i, parse_xml(
                f'<c:dPt {NS}><c:idx val="{i}"/>'
                f'<c:invertIfNegative val="0"/>'
                f'<c:spPr><a:solidFill><a:srgbClr val="{col}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></c:spPr></c:dPt>'
            ))

    def _waterfall_labels(self, series, display_values: list[str]):
        """Boxed value label above each waterfall bar, showing the true signed
        delta (e.g. "-3") rather than the internal stacked/positive-magnitude
        plot value — without this, a bridge chart is unreadable at a glance,
        exactly the gap real IB waterfalls never leave."""
        try:
            # Accessing .data_labels first ensures python-pptx has created a
            # valid, schema-complete <c:dLbls> (group-level show* flags etc.)
            # — individual <c:dLbl> overrides must be inserted before all of
            # that, in ascending idx order, per the CT_DLbls schema sequence.
            series.data_labels.show_value = False
            dLbls = series.data_labels._element
            for i, text in enumerate(display_values):
                dlbl = parse_xml(
                    f'<c:dLbl xmlns:c="{C_NS}" xmlns:a="{A_NS}">'
                    f'<c:idx val="{i}"/>'
                    f'<c:tx><c:rich><a:bodyPr/><a:lstStyle/><a:p><a:r>'
                    f'<a:rPr lang="en-US" sz="700" b="1"><a:solidFill><a:srgbClr val="333333"/></a:solidFill>'
                    f'<a:latin typeface="{self.DEFAULT_FONT}"/></a:rPr>'
                    f'<a:t>{text}</a:t></a:r></a:p></c:rich></c:tx>'
                    f'<c:spPr><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
                    f'<a:ln w="{int(Pt(0.75))}"><a:solidFill><a:srgbClr val="444444"/></a:solidFill></a:ln></c:spPr>'
                    f'<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
                    f'<c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>'
                    f'</c:dLbl>'
                )
                dLbls.insert(i, dlbl)
        except Exception: pass

    # ── Horizontal floating bar (pricing ranges) ───────────────────────────────
    def _hbar_float(self, slide, el: dict):
        series_defs = el.get("series", [])
        if not series_defs:
            return
        labels = [_str(s.get("label", "")) for s in series_defs]
        lows   = [_num(s.get("min"), 0) for s in series_defs]
        highs  = [max(_num(s.get("max"), 0) - _num(s.get("min"), 0), 0) for s in series_defs]

        cd = ChartData()
        cd.categories = labels
        cd.add_series("spacer", lows)
        cd.add_series("range",  highs)

        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart
        self._set_chart_default_font(ch)
        ch.has_title = False; ch.has_legend = False

        sp = ch.series[0]
        sp.format.fill.background()
        sp.format.line.fill.background()

        colors = [self._pcolor(c, self.PALETTE["DKG"]) for c in el.get("colors")] if el.get("colors") else \
            [self.PALETTE["DKG"], self.PALETTE["MDG"], self.PALETTE["OLV"], self.PALETTE["TEL"], self.PALETTE["LBL"], self.PALETTE["GRG"]]
        rng = ch.series[1]
        self._color_bar_points(rng, colors, len(series_defs))

        try:
            ch.value_axis.has_major_gridlines = True
            ch.value_axis.major_gridlines.format.line.color.rgb = _rgb("EBEBEB")
            ch.value_axis.major_gridlines.format.line.width = Pt(0.25)
            ch.value_axis.tick_labels.font.size = Pt(6)
            ch.value_axis.tick_labels.font.name = self.DEFAULT_FONT
            ch.value_axis.tick_labels.font.color.rgb = _rgb("999999")
            ch.value_axis.tick_labels.number_format = "#,##0"
            ch.category_axis.has_major_gridlines = False
            ch.category_axis.tick_labels.font.size = Pt(7)
            ch.category_axis.tick_labels.font.name = self.DEFAULT_FONT
            ch.category_axis.tick_labels.font.color.rgb = _rgb("444444")
        except Exception: pass
        try: ch.plot_area.format.line.fill.background()
        except Exception: pass

    # ── Single-series line chart ───────────────────────────────────────────────
    def _line(self, slide, el: dict):
        labels = el.get("labels") or []
        values = [_num(v) for v in (el.get("values") or [])]
        if not labels or not values:
            return
        cd = ChartData()
        cd.categories = [_str(l) for l in labels]
        cd.add_series("", values)
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        self._set_chart_default_font(ch)
        ser = ch.series[0]
        self._set_line_style(ser, self._pcolor(el.get("color"), self.PALETTE["MDG"]), _num(el.get("width"), 1.8))
        self._smooth_all(ch)
        if el.get("data_labels"):
            self._line_point_labels(ser, el, self.PALETTE["MDG"])
        self._style_axes(ch, ymin=el.get("ymin"), ymax=el.get("ymax"),
                         num_fmt=_str(el.get("num_fmt"), "#,##0"),
                         skip=int(_num(el.get("skip"), 6)), csize=5.5)

    def _line_point_labels(self, ser, el: dict, color: str):
        """Boxed value callout above each point (the 'Historical Performance'
        pattern — every year's revenue/margin number in a thin-bordered box)."""
        try:
            dl = ser.data_labels
            dl.show_value = True
            dl.number_format = _str(el.get("num_fmt"), "#,##0")
            dl.number_format_is_linked = False
            dl.font.size = Pt(7); dl.font.bold = True
            dl.font.name = self.DEFAULT_FONT
            dl.font.color.rgb = _rgb("333333")
            dl.position = XL_LABEL_POSITION.ABOVE
            if el.get("box_labels", True):
                self._box_data_labels(ser, border_color=color)
        except Exception: pass

    # ── Multi-series line chart ────────────────────────────────────────────────
    def _line_multi(self, slide, el: dict):
        series_list = el.get("series") or []
        labels = [_str(l) for l in (el.get("labels") or [])]
        if not labels or not series_list:
            return
        cd = ChartData()
        cd.categories = labels
        for s in series_list:
            vals = [_num(v) for v in (s.get("values") or [])]
            if len(vals) < len(labels):
                last = vals[-1] if vals else 0.0
                vals = vals + [last] * (len(labels) - len(vals))
            cd.add_series(_str(s.get("name"), ""), vals)
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False
        self._set_chart_default_font(ch)
        for i, s in enumerate(series_list):
            if i < len(ch.series):
                col = self._pcolor(s.get("color"), self.PALETTE["DKG"])
                self._set_line_style(ch.series[i], col, _num(s.get("width"), 1.4), dashed=bool(s.get("dashed")))
                if s.get("data_labels"):
                    self._line_point_labels(ch.series[i], el, col)
        self._smooth_all(ch)
        self._style_axes(ch, ymin=_num(el.get("ymin"), 0), ymax=el.get("ymax"),
                         num_fmt=_str(el.get("num_fmt"), "#,##0"),
                         skip=int(_num(el.get("skip"), 1)), csize=5)
        self._style_legend(ch)

    # ── Clustered column chart ─────────────────────────────────────────────────
    def _bar(self, slide, el: dict):
        series_list = el.get("series") or []
        labels = [_str(l) for l in (el.get("labels") or [])]
        if not labels or not series_list:
            return
        cd = ChartData()
        cd.categories = labels
        for s in series_list:
            cd.add_series(_str(s.get("name"), ""), [_num(v) for v in (s.get("values") or [])])
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False
        self._set_chart_default_font(ch)

        for i, s in enumerate(series_list):
            if i >= len(ch.series):
                break
            ser = ch.series[i]
            col = self._pcolor(s.get("color"), self.PALETTE["DKG"])
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = _rgb(col)
            ser.format.line.fill.background()
            if s.get("hatched"):
                self._set_hatch(ser, _str(col, self.PALETTE["DKG"]))
            if s.get("data_labels"):
                try:
                    dl = ser.data_labels
                    dl.show_value = True  # Series has no has_data_labels — this is what actually turns labels on
                    dl.number_format = _str(el.get("num_fmt"), "0%")
                    dl.number_format_is_linked = False
                    dl.font.size = Pt(7); dl.font.bold = False
                    dl.font.name = self.DEFAULT_FONT
                    dl.font.color.rgb = _rgb("444444")
                    dl.position = XL_LABEL_POSITION.OUTSIDE_END
                    # Boxed value callouts are the PANDO default (e.g. "1,727", "37%"
                    # in a thin-bordered box above the bar) — opt out with box_labels:false.
                    if el.get("box_labels", True):
                        self._box_data_labels(ser)
                except Exception: pass

        try:
            ch.plot_area.gap_width = int(_num(el.get("gap_width"), 60))
            ch.plot_area.overlap   = int(_num(el.get("overlap"), -10))
        except Exception: pass

        self._style_axes(ch, ymin=_num(el.get("ymin"), 0), ymax=el.get("ymax"),
                         num_fmt=_str(el.get("num_fmt"), "0%"),
                         skip=int(_num(el.get("skip"), 1)), csize=6.5)

        if len(series_list) > 1:
            self._style_legend(ch)
        else:
            ch.has_legend = False

        # pps-delta annotations + category highlight — the "Brand Perception"
        # pattern: a striped 'perception' bar next to a solid 'experience' bar,
        # with a thin underline + "+Npps" spread call-out above each pair, and
        # an optional dashed box drawing attention to one category.
        if el.get("pair_deltas") and len(series_list) == 2:
            plot_x, plot_y, plot_w, _plot_h = self._plot_rect(x, y, w, h)
            n = len(labels)
            cat_w = plot_w / max(n, 1)
            v0 = [_num(v) for v in (series_list[0].get("values") or [])]
            v1 = [_num(v) for v in (series_list[1].get("values") or [])]
            for i in range(n):
                a = v0[i] if i < len(v0) else 0
                b = v1[i] if i < len(v1) else 0
                cx = plot_x + (i + 0.5) * cat_w
                delta = abs(b - a)
                fmt_val = f"{delta * 100:.0f}pps" if _str(el.get("num_fmt"), "0%").endswith("%") else f"{delta:,.0f}"
                self._txt_box(slide, f"+{fmt_val}", cx - cat_w / 2, plot_y - 0.02, cat_w, 0.18,
                              6.5, align="c", fg="333333")
                try:
                    ln = slide.shapes.add_connector(1, Inches(cx - cat_w * 0.32), Inches(plot_y + 0.16),
                                                    Inches(cx + cat_w * 0.32), Inches(plot_y + 0.16))
                    ln.line.color.rgb = _rgb("333333"); ln.line.width = Pt(0.5)
                except Exception: pass
            hi = el.get("highlight_category")
            if isinstance(hi, int) and 0 <= hi < n:
                hx = plot_x + hi * cat_w
                box = slide.shapes.add_shape(1, Inches(hx + cat_w * 0.05), Inches(plot_y - 0.02),
                                             Inches(cat_w * 0.9), Inches(_plot_h * 0.92))
                box.fill.background()
                box.line.color.rgb = _rgb(el.get("highlight_color", self.PALETTE["OLV"]))
                box.line.width = Pt(0.75)
                self._dash(box)

    def _set_hatch(self, series, color: str):
        try:
            ser_el = series._element
            spPr = ser_el.find("{%s}spPr" % C_NS)
            if spPr is None: return
            solidFill = spPr.find("{%s}solidFill" % A_NS)
            if solidFill is not None:
                spPr.remove(solidFill)
            pattFill = parse_xml(
                # NOTE: "lgDnDiag" is NOT a valid ST_PresetPatternVal (a past bug here
                # silently produced a corrupt .pptx that real PowerPoint refused to
                # open, though python-pptx/LibreOffice didn't complain). "ltDnDiag"
                # is the real value for a fine diagonal-stripe hatch.
                f'<a:pattFill xmlns:a="{A_NS}" prst="ltDnDiag">'
                f'<a:fgClr><a:srgbClr val="{color}"/></a:fgClr>'
                f'<a:bgClr><a:srgbClr val="FFFFFF"/></a:bgClr>'
                f'</a:pattFill>'
            )
            ln_el = spPr.find("{%s}ln" % A_NS)
            if ln_el is not None:
                ln_el.addprevious(pattFill)
            else:
                spPr.append(pattFill)
        except Exception: pass

    # ── Table ─────────────────────────────────────────────────────────────────
    _NUM_CELL_RE = re.compile(r'^[\s$€£]*\(?-?[\d.,]+\)?\s*%?\s*[kKmMbB]?\+?\s*$')

    def _wrap_lines(self, text: str, col_w: float, size: float) -> int:
        """Estimate wrapped-line count for text inside a column of width col_w
        (inches) at font size (pt) — same char-width heuristic used elsewhere
        in this file (icon_row title wrap), applied to table cells so row
        height can be computed from what will actually render, not guessed."""
        text = _str(text)
        if not text:
            return 1
        usable = max(col_w - 0.14, 0.35)
        cpl = max(4, usable / (size * 0.0092))
        total = 0
        for line in text.split("\n"):
            total += max(1, -(-len(line) // int(cpl)))  # ceil div
        return total

    def _table_col_widths(self, el: dict, headers: list, rows: list, w: float) -> list:
        """Distribute column width proportionally to each column's actual content
        length (header + widest cell) rather than splitting evenly — a numbers
        column ('12.4%') and a commentary column (a full sentence) should not
        get the same width, which is how columns used to end up either wasting
        space or wrapping/overflowing their declared row height."""
        n_cols = len(headers)
        explicit = el.get("col_widths")
        if explicit and isinstance(explicit, list) and len(explicit) == n_cols:
            widths = [max(_num(cw, 1.0), 0.3) for cw in explicit]
            total = sum(widths) or 1.0
            return [wd * w / total for wd in widths]
        weights = []
        for c in range(n_cols):
            maxlen = len(_str(headers[c]))
            for row in rows:
                vals = list(row) if isinstance(row, (list, tuple)) else []
                if c < len(vals):
                    maxlen = max(maxlen, len(_str(vals[c])))
            weights.append(max(maxlen, 3))
        total_w = sum(weights) or 1
        raw = [max(w * wt / total_w, 0.55) for wt in weights]
        scale = w / sum(raw) if sum(raw) else 1.0
        return [r * scale for r in raw]

    def _table_layout(self, el: dict):
        """Single source of truth for a table's real geometry — column widths,
        header height, and per-row heights all computed from actual content via
        _wrap_lines, so _table_height (used for slide-layout QA) and _table
        (the actual renderer) can never disagree with each other.

        A table with many rows (e.g. a 9-row operating-lever matrix) can need
        more vertical space than exists between its y and the footer — used to
        just render past the bottom of the slide with the extra rows silently
        clipped off-canvas by PowerPoint. Now: shrink the font first, and if
        that's still not enough (font already at its floor), compress every
        row proportionally so the whole table always fits within the slide."""
        headers = el.get("headers") or []
        rows = el.get("rows") or []
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 3.0))
        col_widths = self._table_col_widths(el, headers, rows, w)

        header_h_base = max(_num(el.get("header_h"), 0.32), 0.1)
        row_h_base    = max(_num(el.get("row_h"),    0.28), 0.1)

        def _compute(size):
            header_lines = 1
            for c, htext in enumerate(headers):
                header_lines = max(header_lines, self._wrap_lines(_str(htext), col_widths[c], size))
            header_line_h = (size * 1.25) / 72 + 0.05
            header_h = max(header_h_base, header_line_h * header_lines + 0.08)

            heights = []
            line_h = (size * 1.22) / 72 + 0.02
            for row in rows:
                vals = list(row) if isinstance(row, (list, tuple)) else []
                max_lines = 1
                for c in range(len(headers)):
                    val = vals[c] if c < len(vals) else ""
                    max_lines = max(max_lines, self._wrap_lines(_str(val), col_widths[c], size))
                heights.append(max(row_h_base, line_h * max_lines + 0.06))
            return header_h, heights

        size = _num(el.get("size"), 8)
        header_h, row_heights = _compute(size)
        total_h = header_h + sum(row_heights)

        available_h = max(CANVAS_BOTTOM - y, 0.5)
        min_size = 6.0
        while total_h > available_h and size > min_size:
            size = max(size - 0.5, min_size)
            header_h, row_heights = _compute(size)
            total_h = header_h + sum(row_heights)

        if total_h > available_h:
            # Font is already at its floor — compress rows proportionally
            # rather than let them render past the canvas edge.
            scale = available_h / total_h
            header_h *= scale
            row_heights = [rh * scale for rh in row_heights]
            total_h = available_h

        return x, y, w, col_widths, header_h, row_heights, max(total_h, 0.2), size

    def _table_height(self, el: dict) -> float:
        """The true rendered height of a table, driven by row count and actual
        text wrap, capped to never exceed the printable canvas — see
        _table_layout. Always trust this over the plan's own declared "h"."""
        _x, _y, _w, _cw, _hh, _rh, total_h, _size = self._table_layout(el)
        return total_h

    def _set_cell_borders(self, cell, color: str = "D9DBD4", width_pt: float = 0.5):
        """Thin borders on every side of a cell — native python-pptx tables render
        with no visible gridlines by default, which reads as a flat, undefined
        block of color rather than a table; real PANDO decks always show hairline
        rules between cells (see the D2C matrix / brand comparison tables)."""
        try:
            tcPr = cell._tc.get_or_add_tcPr()
            w_emu = int(Pt(width_pt))
            # ln* elements must precede fill elements in tcPr's schema order —
            # inserting each at an incrementing index (rather than appending)
            # keeps lnL/lnR/lnT/lnB ahead of whatever fill cell.fill.solid()
            # already wrote, regardless of call order.
            for i, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
                existing = tcPr.find(qn(tag))
                if existing is not None:
                    tcPr.remove(existing)
                ln = parse_xml(
                    f'<{tag} xmlns:a="{A_NS}" w="{w_emu}" cap="flat" cmpd="sng" algn="ctr">'
                    f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                    f'<a:prstDash val="solid"/></{tag}>'
                )
                tcPr.insert(i, ln)
        except Exception: pass

    def _table(self, slide, el: dict):
        headers = el.get("headers") or []
        rows    = el.get("rows")    or []
        if not headers:
            return
        n_cols = len(headers)
        n_rows = len(rows) + 1
        x, y, w, col_widths, header_h, row_heights, h, size = self._table_layout(el)

        gframe = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = gframe.table

        tbl_el = table._tbl
        tblPr = tbl_el.find(qn("a:tblPr"))
        if tblPr is not None:
            tblPr.set("firstRow", "0")
            tblPr.set("bandRow", "0")

        for i, cw in enumerate(col_widths[:n_cols]):
            try: table.columns[i].width = Inches(cw)
            except Exception: pass

        def _style_cell(cell, text, bold, fg, bg, align, border_color="D9DBD4"):
            try:
                cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(bg)
                cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
                cell.margin_top = Inches(0.02);  cell.margin_bottom = Inches(0.02)
                cell.vertical_anchor = 3
                tf = cell.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
                r = p.add_run(); r.text = _str(text)
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
                r.font.bold = bold; r.font.color.rgb = _rgb(fg)
                self._set_cell_borders(cell, border_color)
            except Exception: pass

        for c, htext in enumerate(headers):
            _style_cell(table.cell(0, c), htext, True, self.PALETTE["WHT"], self.PALETTE["DKG"],
                        "l" if c == 0 else "c", border_color=self.PALETTE["DKG"])
        try: table.rows[0].height = Inches(header_h)
        except Exception: pass

        zebra           = el.get("zebra", True)
        bold_first_col  = el.get("bold_first_col", False)
        # label_col: the first column is a row-label rail with alternating dark-green
        # (white text) and white cells — the comparison-matrix style from real PANDO
        # decks (e.g. the BRANDS brand/region matrix). Pairs best with zebra=true.
        label_col = bool(el.get("label_col", False))
        for ridx, row in enumerate(rows):
            bg = self.PALETTE["GRG"] if (zebra and ridx % 2 == 1) else self.PALETTE["WHT"]
            row_vals = list(row) if isinstance(row, (list, tuple)) else []
            for c in range(n_cols):
                val = row_vals[c] if c < len(row_vals) else ""
                if label_col and c == 0:
                    dark = ridx % 2 == 1
                    _style_cell(table.cell(ridx + 1, c), val, False,
                                self.PALETTE["WHT"] if dark else self.PALETTE["NKB"],
                                self.PALETTE["DKG"] if dark else self.PALETTE["WHT"], "l")
                else:
                    # Numeric-looking values ('12.4%', '$1,234', '(3.2)') read better
                    # right-aligned like a spreadsheet; free-text commentary stays
                    # left/centered as before.
                    is_num = c != 0 and bool(self._NUM_CELL_RE.match(_str(val)))
                    align = "r" if is_num else ("l" if c == 0 else "c")
                    _style_cell(table.cell(ridx + 1, c), val, bold_first_col and c == 0,
                                self.PALETTE["NKB"], bg, align)
            try: table.rows[ridx + 1].height = Inches(row_heights[ridx])
            except Exception: pass

    # ── Donut chart ────────────────────────────────────────────────────────────
    def _donut(self, slide, el: dict):
        slices = el.get("slices") or []
        if not slices:
            return
        cd = ChartData()
        cd.categories = [_str(s.get("label", "")) for s in slices]
        cd.add_series("", [_num(s.get("value"), 0) for s in slices])
        x, y, w, h = _geom(el, (0.85, 1.78, 5.0, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False
        self._set_chart_default_font(ch)
        ser = ch.series[0]
        NS = f'xmlns:c="{C_NS}" xmlns:a="{A_NS}"'
        cat_el = ser._element.find(qn("c:cat"))
        if cat_el is None:
            cat_el = ser._element.find(qn("c:val"))
        if cat_el is not None:
            idx = list(ser._element).index(cat_el)
            for i, s in enumerate(slices):
                col = self._pcolor(s.get("color"), self.PALETTE["GRG"])
                ser._element.insert(idx + i, parse_xml(
                    f'<c:dPt {NS}><c:idx val="{i}"/>'
                    f'<c:spPr><a:solidFill><a:srgbClr val="{col}"/></a:solidFill>'
                    f'<a:ln><a:noFill/></a:ln></c:spPr></c:dPt>'
                ))
        hole = int(_num(el.get("hole"), 55))
        for dc in ch._element.findall(".//{%s}doughnutChart" % C_NS):
            hs = dc.find("{%s}holeSize" % C_NS)
            if hs is not None: hs.set("val", str(hole))
            else: dc.append(parse_xml(f'<c:holeSize xmlns:c="{C_NS}" val="{hole}"/>'))

        # Per-slice value/percentage labels — without these the ring is unreadable
        # without a legend (real PANDO donuts, e.g. "Client Profile", show a %
        # figure on every slice). Values are assumed to already be percentage
        # points (23.0, not 0.23); pass num_fmt to override for other units.
        if el.get("data_labels", True) and len(slices) > 1:
            try:
                dl = ser.data_labels
                dl.show_value = True
                dl.number_format = _str(el.get("num_fmt"), '0"%"')
                dl.number_format_is_linked = False
                dl.font.size = Pt(_num(el.get("label_size"), 8))
                dl.font.bold = True
                dl.font.name = self.DEFAULT_FONT
                dl.font.color.rgb = _rgb(el.get("label_color", "FFFFFF" if hole < 70 else "444444"))
            except Exception: pass

        # Bottom legend maps slice color -> category (the ring alone can't carry
        # that once labels are just numbers) — suppress with legend:false when a
        # caller supplies its own legend/annotation instead.
        if el.get("legend", True) and len(slices) > 1:
            self._style_legend(ch)
        else:
            ch.has_legend = False

        # Center KPI text inside the hole ("~570k annual customers / +82% younger
        # than 45") — the Leon-deck donut pattern. First line renders bold; use
        # **bold** spans inside other lines for partial emphasis.
        center = el.get("center") or []
        if isinstance(center, str):
            center = [center]
        center = [_str(c) for c in center if _str(c).strip()]
        if center:
            cw = w * 0.52
            chh = min(0.30 * len(center) + 0.1, h * 0.5)
            tb = slide.shapes.add_textbox(Inches(x + (w - cw) / 2), Inches(y + h / 2 - chh / 2),
                                          Inches(cw), Inches(chh))
            tf = tb.text_frame; tf.word_wrap = True
            try: tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception: pass
            for i, line in enumerate(center):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                for part in re.split(r'(\*\*.*?\*\*)', line):
                    if not part:
                        continue
                    r = p.add_run()
                    if part.startswith("**") and part.endswith("**") and len(part) > 4:
                        r.text = part[2:-2]; r.font.bold = True
                    else:
                        r.text = part; r.font.bold = (i == 0)
                    r.font.name = self.DEFAULT_FONT
                    r.font.size = Pt(_num(el.get("center_size"), 9.5) if i == 0 else _num(el.get("center_size"), 9.5) - 1)
                    r.font.italic = (i > 0)
                    r.font.color.rgb = _rgb(self.PALETTE["NKB"])

    # ── Scatter chart ──────────────────────────────────────────────────────────
    def _scatter(self, slide, el: dict):
        points = el.get("points") or []
        if not points:
            return
        cd = XyChartData()
        for pt in points:
            s = cd.add_series(_str(pt.get("label", "")))
            s.add_data_point(_num(pt.get("x")), _num(pt.get("y")))
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart; ch.has_title = False; ch.has_legend = False
        self._set_chart_default_font(ch)
        for i, pt in enumerate(points):
            if i >= len(ch.series):
                break
            try:
                s = ch.series[i]
                s.format.line.fill.background()
                s.marker.format.fill.solid()
                col = self._pcolor(pt.get("color"), self.PALETTE["DKG"])
                s.marker.format.fill.fore_color.rgb = _rgb(col)
                s.marker.format.line.fill.background()
                s.marker.size = int(_num(pt.get("size"), 10))
            except Exception: pass
            # Native XY scatter can't attach arbitrary text to a point, but each
            # point IS its own series (named by pt["label"]) — showing the series
            # name as the data label is the only way to label a scatter point
            # without hand-rolled pixel math from an axis range we can't know
            # precisely (the chart auto-scales its own axes). python-pptx's
            # Series.data_labels property isn't implemented for XySeries (raises
            # AttributeError), so this writes the <c:dLbls> OOXML directly.
            if _str(pt.get("label")):
                self._scatter_point_label(s, col, _num(pt.get("label_size"), 7))
        try:
            ch.value_axis.tick_labels.font.size = Pt(6)
            ch.value_axis.tick_labels.font.name = self.DEFAULT_FONT
            ch.value_axis.tick_labels.number_format = _str(el.get("y_fmt"), "#,##0")
            ch.value_axis.tick_labels.number_format_is_linked = False
            ch.category_axis.tick_labels.font.size = Pt(6)
            ch.category_axis.tick_labels.font.name = self.DEFAULT_FONT
            ch.category_axis.tick_labels.number_format = _str(el.get("x_fmt"), "#,##0")
            ch.category_axis.tick_labels.number_format_is_linked = False
        except Exception: pass

    def _scatter_point_label(self, series, color: str, size: float = 7):
        """Insert a <c:dLbls> block showing the series name, positioned to the
        right of the marker — the only way to label an XY scatter point since
        python-pptx exposes no data_labels API for XySeries."""
        try:
            ser_el = series._element
            dLbls = parse_xml(
                f'<c:dLbls xmlns:c="{C_NS}" xmlns:a="{A_NS}">'
                f'<c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'
                f'<c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>'
                f'<a:defRPr sz="{int(size * 100)}" b="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'<a:latin typeface="{self.DEFAULT_FONT}"/></a:defRPr>'
                f'</a:pPr><a:endParaRPr lang="en-US"/></a:p></c:txPr>'
                f'<c:dLblPos val="r"/>'
                f'<c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>'
                f'<c:showSerName val="1"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>'
                f'</c:dLbls>'
            )
            xVal = ser_el.find(qn("c:xVal"))
            if xVal is not None:
                xVal.addprevious(dLbls)
            else:
                ser_el.append(dLbls)
        except Exception: pass

    # ── Quadrant map ───────────────────────────────────────────────────────────
    def _quadrant(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        labels = el.get("axis_labels") or {}
        mid_x = x + w / 2; mid_y = y + h / 2
        try:
            vl = slide.shapes.add_connector(1, Inches(mid_x), Inches(y), Inches(mid_x), Inches(y + h))
            vl.line.color.rgb = _rgb("CCCCCC"); vl.line.width = Pt(0.5)
            hl = slide.shapes.add_connector(1, Inches(x), Inches(mid_y), Inches(x + w), Inches(mid_y))
            hl.line.color.rgb = _rgb("CCCCCC"); hl.line.width = Pt(0.5)
        except Exception: pass

        def _lbl(text, lx, ly, lw=1.5, lh=0.22, size=6.5, fg="888888", align="c"):
            try:
                box = slide.shapes.add_textbox(Inches(lx), Inches(ly), Inches(lw), Inches(lh))
                tf = box.text_frame; p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if align == "c" else PP_ALIGN.LEFT
                r = p.add_run(); r.text = _str(text)
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
                r.font.color.rgb = _rgb(fg)
            except Exception: pass

        # Clamp left/right axis labels inside the printable canvas — a quadrant
        # sized close to the full content width used to push the right-side
        # label (placed at x+w+0.05) past CONTENT_RIGHT and off the slide.
        right_lx = min(x + w + 0.05, CONTENT_RIGHT - 1.5)
        left_lx  = max(x - 1.6, CONTENT_X - 0.05)
        if labels.get("top"):    _lbl(labels["top"],    mid_x - 0.75, y - 0.25)
        if labels.get("bottom"): _lbl(labels["bottom"], mid_x - 0.75, y + h + 0.03)
        if labels.get("left"):   _lbl(labels["left"],   left_lx,      mid_y - 0.10)
        if labels.get("right"):  _lbl(labels["right"],  right_lx,     mid_y - 0.10)

        # Two dots placed close together end up with their labels overlapping
        # (both default to sitting just below the dot) — alternate a nearby
        # label above the dot instead, a cheap stand-in for real collision
        # detection that catches the common "two competitors near each other"
        # case seen in the reference deck's positioning maps.
        placed: list[tuple[float, float]] = []
        for brand in (el.get("brands") or []):
            try:
                bx = x + _num(brand.get("px")) * w
                by = y + (1 - _num(brand.get("py"))) * h
                col = self._pcolor(brand.get("color"), self.PALETTE["DKG"])
                dot = slide.shapes.add_shape(9, Inches(bx - 0.08), Inches(by - 0.08), Inches(0.16), Inches(0.16))
                dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
                dot.line.fill.background()
                below = not any(abs(bx - px) < 0.9 and abs(by - py) < 0.35 for px, py in placed)
                placed.append((bx, by))
                label_y = by + 0.10 if below else by - 0.30
                box = slide.shapes.add_textbox(Inches(bx - 0.5), Inches(label_y), Inches(1.0), Inches(0.20))
                tf = box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = _str(brand.get("label"))
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(6.5)
                r.font.color.rgb = _rgb(col)
            except Exception: pass

    # ── Band scatter: shaded range bands + dot plot + brand legend ─────────────
    # The "Store Payback" pattern — points positioned by fractional px/py (same
    # convention as _quadrant) over shaded horizontal bands with boxed bucket-%
    # labels on the left, plus a bottom row of colored-dot brand legend.
    def _band_scatter(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)

        legend = [g for g in (el.get("legend") or []) if isinstance(g, dict)]
        legend_h = 0.22 if legend else 0.0
        x_ticks = el.get("x_ticks") or []
        xtick_h = 0.20 if x_ticks else 0.0
        plot_x, plot_w = x + 0.42, w - 0.42 - 0.08
        plot_y, plot_h = y, h - legend_h - xtick_h - 0.03

        for band in (el.get("bands") or []):
            y0 = max(0.0, min(1.0, _num(band.get("y0"), 0)))
            y1 = max(0.0, min(1.0, _num(band.get("y1"), 1)))
            by = plot_y + (1 - y1) * plot_h
            bh = max((y1 - y0) * plot_h, 0.01)
            rect = slide.shapes.add_shape(1, Inches(plot_x), Inches(by), Inches(plot_w), Inches(bh))
            rect.fill.solid(); rect.fill.fore_color.rgb = _rgb(self._pcolor(band.get("color"), "F2F2F2"))
            rect.line.fill.background()
            lbl = _str(band.get("label"))
            if lbl and bh >= 0.18:  # skip label on a band too thin to hold a 0.22"-tall box without overlapping neighbors
                tb = slide.shapes.add_shape(1, Inches(plot_x + 0.06), Inches(by + bh / 2 - 0.11),
                                            Inches(1.15), Inches(0.22))
                tb.fill.solid(); tb.fill.fore_color.rgb = _rgb("FFFFFF")
                tb.line.color.rgb = _rgb("999999"); tb.line.width = Pt(0.5)
                tf = tb.text_frame; tf.word_wrap = False
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = lbl
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(6.5); r.font.color.rgb = _rgb("444444")

        try:
            ax = slide.shapes.add_connector(1, Inches(plot_x), Inches(plot_y), Inches(plot_x), Inches(plot_y + plot_h))
            ax.line.color.rgb = _rgb("CCCCCC"); ax.line.width = Pt(0.75)
            axx = slide.shapes.add_connector(1, Inches(plot_x), Inches(plot_y + plot_h),
                                             Inches(plot_x + plot_w), Inches(plot_y + plot_h))
            axx.line.color.rgb = _rgb("CCCCCC"); axx.line.width = Pt(0.75)
        except Exception: pass

        for t in (el.get("y_ticks") or []):
            py = max(0.0, min(1.0, _num(t.get("py"), 0)))
            ty = plot_y + (1 - py) * plot_h
            self._txt_box(slide, _str(t.get("label")), x, ty - 0.08, 0.36, 0.16, 6, align="r", fg="999999")

        n = len(x_ticks)
        if n:
            ty = plot_y + plot_h + 0.03
            for i, lbl in enumerate(x_ticks):
                px = i / (n - 1) if n > 1 else 0.5
                tx = plot_x + px * plot_w
                self._txt_box(slide, _str(lbl), tx - 0.4, ty, 0.8, 0.18, 6.5, align="c", fg="999999")

        for pt in (el.get("points") or []):
            px = max(0.0, min(1.0, _num(pt.get("px"))))
            py = max(0.0, min(1.0, _num(pt.get("py"))))
            cx = plot_x + px * plot_w
            cy = plot_y + (1 - py) * plot_h
            d = _num(pt.get("size"), 0.09)
            dot = slide.shapes.add_shape(9, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(self._pcolor(pt.get("color"), self.PALETTE["DKG"]))
            dot.line.fill.background()

        if legend:
            slot_w = w / len(legend)
            ly = y + h - legend_h
            for i, g in enumerate(legend):
                lx = x + i * slot_w
                dot = slide.shapes.add_shape(9, Inches(lx), Inches(ly + 0.03), Inches(0.09), Inches(0.09))
                dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(self._pcolor(g.get("color"), self.PALETTE["DKG"]))
                dot.line.fill.background()
                self._txt_box(slide, _str(g.get("label")), lx + 0.14, ly, slot_w - 0.16, 0.18, 7, fg="444444")

    # ── Shared text/bullet helpers for the new content elements ────────────────
    def _txt_box(self, slide, text, x, y, w, h, size, bold=False, italic=False,
                 fg="0A231F", align="l", anchor=None, wrap=True):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = wrap
        if anchor is not None:
            try: tf.vertical_anchor = anchor
            except Exception: pass
        align_map = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
        lines = _str(text).split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            # **bold** spans render as bold runs — the Leon-deck pattern of body
            # copy with only the key figures/phrases emphasized.
            for part in re.split(r'(\*\*.*?\*\*)', line):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    r.text = part[2:-2]; r.font.bold = True
                else:
                    r.text = part; r.font.bold = bold
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
                r.font.italic = italic
                r.font.color.rgb = _rgb(fg)
        return box

    def _dash(self, line_or_conn, style: str = "dash"):
        """Set a dashed stroke on a connector/shape line (grey QA-friendly dashes)."""
        try:
            ln = line_or_conn.line._get_or_add_ln()
            pd = ln.find("{%s}prstDash" % A_NS)
            if pd is None:
                pd = etree.SubElement(ln, "{%s}prstDash" % A_NS)
            pd.set("val", style)
        except Exception: pass

    def _chart_header(self, slide, el: dict, x: float, y: float, w: float) -> float:
        """Plain-text header above a chart — bold title + italic grey subtitle
        (the 'Brand NPS / (Net Promoter Score)' pattern from real PANDO decks).
        Returns the vertical space consumed so the chart starts below it."""
        title = _str(el.get("title")).strip()
        sub   = _str(el.get("subtitle")).strip()
        dy = 0.0
        if title:
            self._txt_box(slide, title, x, y, w, 0.24, 10.5, bold=True, fg=self.PALETTE["NKB"])
            dy += 0.26
        if sub:
            self._txt_box(slide, sub, x, y + dy, w, 0.20, 8.5, italic=True, fg="666666")
            dy += 0.24
        return dy

    def _bullet_list(self, slide, items, x, y, w, h, size=8.5, fg="0A231F", bullet_color=None, space_after_pt=3):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        bcol = bullet_color or self.PALETTE["DKG"]
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(space_after_pt)
            pPr = p._p.get_or_add_pPr()
            pPr.set("indent", "-137160"); pPr.set("marL", "137160")
            buFont = parse_xml(f'<a:buFont xmlns:a="{A_NS}" typeface="Arial"/>')
            buChar = parse_xml(f'<a:buChar xmlns:a="{A_NS}" char="&#8226;"/>')
            buClr = parse_xml(f'<a:buClr xmlns:a="{A_NS}"><a:srgbClr val="{bcol}"/></a:buClr>')
            pPr.append(buClr); pPr.append(buFont); pPr.append(buChar)
            for part in re.split(r'(\*\*.*?\*\*)', _str(item)):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    r.text = part[2:-2]; r.font.bold = True
                else:
                    r.text = part
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(size)
                r.font.color.rgb = _rgb(fg)
        return box

    # ── Stat row: large-number KPI callouts (no icons, just number + label) ────
    def _stat_row(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 1.6))
        items = [it for it in (el.get("items") or []) if isinstance(it, dict)]
        if not items:
            return
        n = len(items)
        gap = 0.35
        cell_w = (w - gap * (n - 1)) / n
        value_size = _num(el.get("value_size"), 32)
        label_size = _num(el.get("label_size"), 9)
        for i, it in enumerate(items):
            cx = x + i * (cell_w + gap)
            col = self._pcolor(it.get("color"), self.PALETTE["DKG"])
            sub = _str(it.get("delta"))
            # Reserve a top band for the delta badge so a long/wrapping value never
            # collides with it — they used to share the same y and overlapped
            # whenever the value text wrapped to two lines in a narrow cell.
            delta_h = 0.22 if sub else 0.0
            if sub:
                self._txt_box(slide, sub, cx, y, cell_w, delta_h,
                              label_size - 0.5, italic=True, fg=col, align="r")
            # Long values in a crowded row (5+ items) need a smaller size to have
            # any chance of staying on one line instead of wrapping into the label.
            v_size = value_size if n <= 4 else min(value_size, 22)
            self._txt_box(slide, _str(it.get("value")), cx, y + delta_h, cell_w, h * 0.62 - delta_h,
                          v_size, bold=True, fg=col, align="l")
            label_y = y + h * 0.62
            self._txt_box(slide, _str(it.get("label")), cx, label_y, cell_w, h * 0.30,
                          label_size, fg="555555", align="l")
            if i < n - 1:
                try:
                    ln = slide.shapes.add_connector(1, Inches(cx + cell_w + gap / 2), Inches(y + 0.05),
                                                    Inches(cx + cell_w + gap / 2), Inches(y + h - 0.05))
                    ln.line.color.rgb = _rgb(self.PALETTE["GRG"]); ln.line.width = Pt(0.5)
                except Exception: pass

    # ── Icon row: colored glyph circle + bold header + description ─────────────
    def _icon_row(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 3.5))
        items = [it for it in (el.get("items") or []) if isinstance(it, dict)]
        if not items:
            return
        direction = _str(el.get("direction"), "col")
        n = len(items)
        circle_d = _num(el.get("circle_size"), 0.42)

        def _one(ix, iy, iw, ih, it):
            col = self._pcolor(it.get("color"), self.PALETTE["DKG"])
            dot = slide.shapes.add_shape(9, Inches(ix), Inches(iy), Inches(circle_d), Inches(circle_d))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
            dot.line.fill.background()
            tf = dot.text_frame; tf.word_wrap = False
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = _str(it.get("glyph"), "•")
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(_num(it.get("glyph_size"), 13))
            r.font.bold = True; r.font.color.rgb = _rgb(self.PALETTE["WHT"])
            text_x = ix + circle_d + 0.15
            text_w = iw - circle_d - 0.15
            title = _str(it.get("title"))
            title_size = _num(it.get("title_size"), 10.5)
            # A narrow column (e.g. 4+ items in a row) often wraps the title to two
            # lines — reserve a second line of height so the body copy doesn't start
            # underneath the wrapped second line of the title.
            chars_per_line = max(1, text_w / (title_size * 0.0095))
            title_lines = 2 if len(title) > chars_per_line else 1
            title_h = 0.24 * title_lines
            self._txt_box(slide, title, text_x, iy - 0.03, text_w, title_h,
                          title_size, bold=True, fg=self.PALETTE["NKB"])
            body = _str(it.get("text"))
            if body:
                self._txt_box(slide, body, text_x, iy - 0.03 + title_h, text_w, ih - title_h,
                              _num(it.get("text_size"), 8), fg="555555")

        if direction == "row":
            # A single row only has room for a handful of items before the text
            # column becomes too narrow to hold a title, let alone body copy —
            # wrap into a multi-row grid once items would drop below a readable
            # minimum width instead of squeezing everything onto one line.
            gap = 0.3
            min_text_w = 1.3
            min_cell_w = circle_d + 0.15 + min_text_w
            cols = max(1, min(n, int((w + gap) // (min_cell_w + gap))))
            rows = -(-n // cols)  # ceil
            cell_w = (w - gap * (cols - 1)) / cols
            cell_h = (h - gap * (rows - 1)) / rows
            for i, it in enumerate(items):
                r, c = divmod(i, cols)
                _one(x + c * (cell_w + gap), y + r * (cell_h + gap), cell_w, cell_h, it)
        else:
            row_h = h / n
            for i, it in enumerate(items):
                _one(x, y + i * row_h, w, row_h - 0.08, it)

    # ── Comparison cards: 2-4 subtly-tinted panels with a header + bullets ─────
    def _comparison_cards(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        cards = [c for c in (el.get("cards") or []) if isinstance(c, dict)]
        if not cards:
            return
        n = len(cards)
        gap = 0.28
        card_w = (w - gap * (n - 1)) / n
        pad = 0.18
        for i, card in enumerate(cards):
            cx = x + i * (card_w + gap)
            col = self._pcolor(card.get("color"), self.PALETTE["DKG"])
            tint = _tint(col, 0.90)
            panel = slide.shapes.add_shape(1, Inches(cx), Inches(y), Inches(card_w), Inches(h))
            panel.fill.solid(); panel.fill.fore_color.rgb = _rgb(tint)
            panel.line.color.rgb = _rgb(self.PALETTE["GRG"]); panel.line.width = Pt(0.5)
            self._txt_box(slide, _str(card.get("title")), cx + pad, y + 0.14, card_w - pad * 2, 0.32,
                          _num(card.get("title_size"), 11), bold=True, fg=col)
            bullets = card.get("bullets") or []
            if bullets:
                self._bullet_list(slide, bullets, cx + pad, y + 0.55, card_w - pad * 2, h - 0.7,
                                  size=_num(card.get("text_size"), 8.5), fg=self.PALETTE["NKB"], bullet_color=col)

    # ── Timeline: numbered milestones connected by a horizontal line ───────────
    def _timeline(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 3.0))
        steps = [s for s in (el.get("steps") or []) if isinstance(s, dict)]
        n = len(steps)
        if not n:
            return
        dot_d = _num(el.get("dot_size"), 0.32)
        line_y = y + dot_d / 2
        try:
            ln = slide.shapes.add_connector(1, Inches(x + dot_d / 2), Inches(line_y),
                                            Inches(x + w - dot_d / 2), Inches(line_y))
            ln.line.color.rgb = _rgb(self.PALETTE["GRG"]); ln.line.width = Pt(1.5)
        except Exception: pass

        cycle = [self.PALETTE["DKG"], self.PALETTE["MDG"], self.PALETTE["OLV"],
                 self.PALETTE["TEL"], self.PALETTE["LBL"]]
        slot_w = w / n if n > 1 else w
        for i, step in enumerate(steps):
            cx = x + slot_w * i + (slot_w - dot_d) / 2 if n > 1 else x
            col = self._pcolor(step.get("color"), cycle[i % len(cycle)])
            dot = slide.shapes.add_shape(9, Inches(cx), Inches(y), Inches(dot_d), Inches(dot_d))
            dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
            dot.line.color.rgb = _rgb(self.PALETTE["WHT"]); dot.line.width = Pt(1.25)
            tf = dot.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = _str(step.get("num"), str(i + 1))
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(_num(step.get("num_size"), 10))
            r.font.bold = True; r.font.color.rgb = _rgb(self.PALETTE["WHT"])

            label_w = slot_w - 0.15
            label_x = x + slot_w * i + 0.075 if n > 1 else x
            self._txt_box(slide, _str(step.get("label")), label_x, y + dot_d + 0.10, label_w, 0.24,
                          _num(step.get("label_size"), 9), bold=True, align="c", fg=self.PALETTE["NKB"])
            body = _str(step.get("text"))
            if body:
                self._txt_box(slide, body, label_x, y + dot_d + 0.36, label_w, h - dot_d - 0.36,
                              _num(step.get("text_size"), 7.5), align="c", fg="666666")

    # ── Waterfall / bridge chart (e.g. Revenue → EBITDA walk) ───────────────────
    def _waterfall(self, slide, el: dict):
        labels = [_str(l) for l in (el.get("labels") or [])]
        values = [_num(v) for v in (el.get("values") or [])]
        if not labels or not values or len(labels) != len(values):
            return
        totals_flags = el.get("totals") or []
        is_total = [bool(totals_flags[i]) if i < len(totals_flags) else (i == 0 or i == len(values) - 1)
                    for i in range(len(values))]

        bases, tops, colors = [], [], []
        up_color    = _str(el.get("up_color"),    self.PALETTE["MDG"])
        down_color  = _str(el.get("down_color"),  self.PALETTE["TEL"])
        total_color = _str(el.get("total_color"), self.PALETTE["DKG"])
        cum = 0.0
        for v, tot in zip(values, is_total):
            if tot:
                bases.append(0.0); tops.append(v); colors.append(total_color)
                cum = v
            elif v >= 0:
                bases.append(cum); tops.append(v); colors.append(up_color)
                cum += v
            else:
                cum += v
                bases.append(cum); tops.append(-v); colors.append(down_color)

        cd = ChartData()
        cd.categories = labels
        cd.add_series("base", bases)
        cd.add_series("delta", tops)

        x, y, w, h = _geom(el, (0.85, 1.78, 12.18, 4.0))
        hdr = self._chart_header(slide, el, x, y, w)
        y += hdr; h = max(h - hdr, 0.8)
        cf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        ch = cf.chart
        ch.has_title = False; ch.has_legend = False
        self._set_chart_default_font(ch)

        base_series = ch.series[0]
        base_series.format.fill.background()
        base_series.format.line.fill.background()

        delta_series = ch.series[1]
        self._color_bar_points(delta_series, colors, len(values))
        if el.get("data_labels", True):
            display = [f"{v:,.0f}" if tot else f"{'+' if v >= 0 else '-'}{abs(v):,.0f}"
                       for v, tot in zip(values, is_total)]
            self._waterfall_labels(delta_series, display)

        gap_width = int(_num(el.get("gap_width"), 45))
        try:
            ch.plot_area.gap_width = gap_width
        except Exception: pass

        # Force an explicit axis max (rather than leaving it to PowerPoint's
        # auto-scale) so the bridge-line overlay below — drawn with our own
        # approximate plot-rect math — lines up with where the bars actually
        # render instead of guessing a scale PowerPoint might pick differently.
        top_edges = [b + t for b, t in zip(bases, tops)]
        ymax_eff = _num(el.get("ymax")) or (max(top_edges) * 1.15 if top_edges else 1.0)
        self._style_axes(ch, ymin=_num(el.get("ymin"), 0), ymax=ymax_eff,
                         num_fmt=_str(el.get("num_fmt"), "#,##0"),
                         skip=1, csize=6.5)

        # Dashed bridge lines connecting the top of each bar to the next — without
        # these, a stacked-column waterfall is visually indistinguishable from a
        # plain stacked bar chart (the defining "waterfall" cue is the connector).
        if el.get("bridges", True) and len(values) > 1:
            # top=0.16 (vs. _plot_rect's 0.05 default) calibrated against a
            # COM-rendered waterfall: a single-series stacked column with no
            # legend reserves more top headroom than the bar-chart pair_deltas
            # overlay this helper was originally tuned for.
            plot_x, plot_y, plot_w, plot_h = self._plot_rect(x, y, w, h, top=0.16)
            n = len(values)
            cat_w = plot_w / n
            bar_frac = 100.0 / (100.0 + gap_width)
            for i in range(n - 1):
                edge_val = top_edges[i]
                ty = plot_y + plot_h * (1 - (edge_val / ymax_eff if ymax_eff else 0))
                x1 = plot_x + (i + 0.5) * cat_w + cat_w * bar_frac / 2
                x2 = plot_x + (i + 1.5) * cat_w - cat_w * bar_frac / 2
                try:
                    conn = slide.shapes.add_connector(1, Inches(x1), Inches(ty), Inches(x2), Inches(ty))
                    conn.line.color.rgb = _rgb("999999"); conn.line.width = Pt(0.75)
                    self._dash(conn)
                except Exception: pass

    # ── Alternating timeline: entries above/below a central axis ───────────────
    # The COMPANY HISTORY pattern: a horizontal line mid-element, milestone dots on
    # it, entries alternating above (even index) and below (odd index) with dashed
    # connectors, each entry a bold accent-colored year plus **bold**-mixed text.
    def _alt_timeline(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 2.0, 12.18, 4.4))
        entries = [e for e in (el.get("entries") or []) if isinstance(e, dict)]
        n = len(entries)
        if not n:
            return
        line_y = y + h * 0.5
        dot_d = 0.12
        try:
            ln = slide.shapes.add_connector(1, Inches(x), Inches(line_y), Inches(x + w), Inches(line_y))
            ln.line.color.rgb = _rgb(self.PALETTE["NKB"]); ln.line.width = Pt(1.25)
        except Exception: pass

        slot_w = w / n
        text_h = h * 0.5 - 0.35
        for i, e in enumerate(entries):
            cx = x + slot_w * i + slot_w / 2
            col = self._pcolor(e.get("color"), self.PALETTE["DKG"])
            above = i % 2 == 0
            try:
                dot = slide.shapes.add_shape(9, Inches(cx - dot_d / 2), Inches(line_y - dot_d / 2),
                                             Inches(dot_d), Inches(dot_d))
                dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(col)
                dot.line.fill.background()
            except Exception: pass
            conn_len = 0.30
            try:
                if above:
                    c1 = slide.shapes.add_connector(1, Inches(cx), Inches(line_y - dot_d / 2 - conn_len),
                                                    Inches(cx), Inches(line_y - dot_d / 2))
                else:
                    c1 = slide.shapes.add_connector(1, Inches(cx), Inches(line_y + dot_d / 2),
                                                    Inches(cx), Inches(line_y + dot_d / 2 + conn_len))
                c1.line.color.rgb = _rgb("999999"); c1.line.width = Pt(0.75)
                self._dash(c1)
            except Exception: pass

            tx = x + slot_w * i + 0.06
            tw = slot_w - 0.12
            if above:
                tb = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(tw), Inches(text_h))
                anchor = MSO_ANCHOR.BOTTOM
            else:
                tb = slide.shapes.add_textbox(Inches(tx), Inches(line_y + dot_d / 2 + conn_len + 0.05),
                                              Inches(tw), Inches(text_h))
                anchor = MSO_ANCHOR.TOP
            tf = tb.text_frame; tf.word_wrap = True
            try: tf.vertical_anchor = anchor
            except Exception: pass
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = _str(e.get("label"))
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(_num(e.get("label_size"), 10))
            r.font.bold = True; r.font.color.rgb = _rgb(col)
            body = _str(e.get("text"))
            if body:
                p2 = tf.add_paragraph()
                for part in re.split(r'(\*\*.*?\*\*)', body):
                    if not part:
                        continue
                    r2 = p2.add_run()
                    if part.startswith("**") and part.endswith("**") and len(part) > 4:
                        r2.text = part[2:-2]; r2.font.bold = True
                    else:
                        r2.text = part
                    r2.font.name = self.DEFAULT_FONT
                    r2.font.size = Pt(_num(e.get("text_size"), 7.5))
                    r2.font.color.rgb = _rgb("444444")

    # ── Org chart: hierarchy of two-band boxes with dashed connectors ──────────
    # The CORPORATE STRUCTURE pattern: each node is a colored title band over a
    # light subtitle band (jurisdiction), an optional italic note below, and a
    # dashed connector from its parent annotated with an ownership percentage.
    def _org_chart(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 2.0, 12.18, 4.4))
        levels = [lv for lv in (el.get("levels") or []) if isinstance(lv, list) and lv]
        n_levels = len(levels)
        if not n_levels:
            return
        title_h = 0.34; sub_h = 0.26
        level_h = h / n_levels
        positions: list[list[tuple[float, float]]] = []  # (center_x, box_bottom_y) per node

        cycle = [self.PALETTE["NKB"], self.PALETTE["MDG"], self.PALETTE["OLV"]]
        for li, level in enumerate(levels):
            nodes = [nd for nd in level if isinstance(nd, dict)]
            k = len(nodes)
            slot_w = w / k
            box_w = min(slot_w - 0.4, 3.6)
            ly = y + li * level_h
            row_pos = []
            for ni, nd in enumerate(nodes):
                cx = x + slot_w * ni + slot_w / 2
                bx = cx - box_w / 2
                col = self._pcolor(nd.get("color"), cycle[li % len(cycle)])
                bar = slide.shapes.add_shape(1, Inches(bx), Inches(ly), Inches(box_w), Inches(title_h))
                bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(col)
                bar.line.fill.background()
                tf = bar.text_frame; tf.word_wrap = False
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = _str(nd.get("label"))
                r.font.name = self.DEFAULT_FONT; r.font.size = Pt(9.5)
                r.font.bold = False; r.font.color.rgb = _rgb(self.PALETTE["WHT"])
                sub = _str(nd.get("sub"))
                bottom = ly + title_h
                if sub:
                    sb = slide.shapes.add_shape(1, Inches(bx), Inches(bottom), Inches(box_w), Inches(sub_h))
                    sb.fill.solid(); sb.fill.fore_color.rgb = _rgb(_tint(self.PALETTE["GRG"], 0.4))
                    sb.line.fill.background()
                    tf2 = sb.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
                    r2 = p2.add_run(); r2.text = sub
                    r2.font.name = self.DEFAULT_FONT; r2.font.size = Pt(8)
                    r2.font.italic = True; r2.font.color.rgb = _rgb("555555")
                    bottom += sub_h
                note = _str(nd.get("note"))
                if note:
                    self._txt_box(slide, note, bx, bottom + 0.04, box_w, 0.22, 8,
                                  italic=True, fg="777777", align="c")
                row_pos.append((cx, bottom))

                parent_idx = nd.get("parent")
                pct = _str(nd.get("pct"))
                if li > 0 and parent_idx is not None and int(_num(parent_idx, 0)) < len(positions[li - 1]):
                    px, pbottom = positions[li - 1][int(_num(parent_idx, 0))]
                    try:
                        conn = slide.shapes.add_connector(2, Inches(px), Inches(pbottom),
                                                          Inches(cx), Inches(ly))
                        conn.line.color.rgb = _rgb("999999"); conn.line.width = Pt(0.75)
                        self._dash(conn)
                    except Exception: pass
                    if pct:
                        self._txt_box(slide, pct, cx + 0.08, ly - 0.30, 0.9, 0.22, 8,
                                      italic=True, fg="444444")
            positions.append(row_pos)

    # ── Process flow: boxed steps connected by arrows ───────────────────────────
    # For value chains and step sequences (Designer → Manufacturer → ... → Retailer):
    # each step is a dark header band with description text below, joined by arrows.
    def _process_flow(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 2.0, 12.18, 2.2))
        steps = [s for s in (el.get("steps") or []) if isinstance(s, dict)]
        n = len(steps)
        if not n:
            return
        arrow_w = 0.32
        gap = 0.12
        box_w = (w - (arrow_w + gap * 2) * (n - 1)) / n
        head_h = 0.34
        for i, s in enumerate(steps):
            bx = x + i * (box_w + arrow_w + gap * 2)
            col = self._pcolor(s.get("color"), self.PALETTE["NKB"])
            bar = slide.shapes.add_shape(1, Inches(bx), Inches(y), Inches(box_w), Inches(head_h))
            bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(col)
            bar.line.fill.background()
            tf = bar.text_frame; tf.word_wrap = False
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = _str(s.get("title"))
            r.font.name = self.DEFAULT_FONT; r.font.size = Pt(_num(s.get("title_size"), 9.5))
            r.font.color.rgb = _rgb(self.PALETTE["WHT"])
            body = _str(s.get("text"))
            if body:
                self._txt_box(slide, body, bx + 0.05, y + head_h + 0.08, box_w - 0.1,
                              h - head_h - 0.08, _num(s.get("text_size"), 7.5),
                              italic=True, fg="444444", align="c")
            if i < n - 1:
                try:
                    ar = slide.shapes.add_shape(33, Inches(bx + box_w + gap),
                                                Inches(y + head_h / 2 - 0.07),
                                                Inches(arrow_w), Inches(0.14))
                    ar.fill.solid(); ar.fill.fore_color.rgb = _rgb(self.PALETTE["NKB"])
                    ar.line.fill.background()
                except Exception: pass

    # ── Pill row: rounded-rectangle stat callouts ───────────────────────────────
    # The bottom-of-slide KPI band ("+3.0x sales per store vs. incumbents"): tinted
    # rounded rectangles with **bold**-mixed sentences. Distinct from stat_row —
    # these are sentences with an embedded figure, not a big number over a label.
    def _pill_row(self, slide, el: dict):
        x, y, w, h = _geom(el, (0.85, 5.6, 12.18, 0.75))
        items = [it for it in (el.get("items") or []) if isinstance(it, dict) or isinstance(it, str)]
        n = len(items)
        if not n:
            return
        gap = 0.28
        pill_w = (w - gap * (n - 1)) / n
        for i, it in enumerate(items):
            if isinstance(it, str):
                it = {"text": it}
            col = self._pcolor(it.get("color"), self.PALETTE["MDG"])
            px = x + i * (pill_w + gap)
            pill = slide.shapes.add_shape(5, Inches(px), Inches(y), Inches(pill_w), Inches(h))
            pill.fill.solid(); pill.fill.fore_color.rgb = _rgb(_tint(col, 0.82))
            pill.line.fill.background()
            tf = pill.text_frame; tf.word_wrap = True
            tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            for part in re.split(r'(\*\*.*?\*\*)', _str(it.get("text"))):
                if not part:
                    continue
                r = p.add_run()
                if part.startswith("**") and part.endswith("**") and len(part) > 4:
                    r.text = part[2:-2]; r.font.bold = True
                else:
                    r.text = part
                r.font.name = self.DEFAULT_FONT
                r.font.size = Pt(_num(it.get("size"), 9))
                r.font.color.rgb = _rgb(self.PALETTE["NKB"])
