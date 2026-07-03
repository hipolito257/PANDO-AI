"""
TemplateProfiler — extracts colors, fonts, layouts from a PPTX template.
Returns a profile dict the Next.js app stores in the DB and sends to Claude.
"""
import io, re, zipfile
from xml.etree import ElementTree as ET

from lxml import etree
from pptx import Presentation
from pptx.util import Inches


THEME_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PH_NAMES = {18: "category", 16: "title", 17: "takeaway", 14: "note", 26: "content"}


def _luminance(hex_color: str) -> float:
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _color_distance(hex_a: str, hex_b: str) -> float:
    ra, ga, ba = int(hex_a[0:2], 16), int(hex_a[2:4], 16), int(hex_a[4:6], 16)
    rb, gb, bb = int(hex_b[0:2], 16), int(hex_b[2:4], 16), int(hex_b[4:6], 16)
    return ((ra - rb) ** 2 + (ga - gb) ** 2 + (ba - bb) ** 2) ** 0.5


class TemplateProfiler:
    def __init__(self, template_bytes: bytes):
        self.prs = Presentation(io.BytesIO(template_bytes))
        self.raw = template_bytes

    def extract(self) -> dict:
        colors = self._colors()
        return {
            "colors":  colors,
            "palette": self._map_palette(colors),
            "fonts":   self._fonts(),
            "layouts": self._layouts(),
            "slide_width_in":  round(self.prs.slide_width  / 914400, 2),
            "slide_height_in": round(self.prs.slide_height / 914400, 2),
        }

    def _map_palette(self, colors: dict) -> dict:
        """Map the template's real accent colors onto the semantic palette keys
        PptxBuilder actually draws with (DKG/MDG/OLV/TEL/LBL/GRG/NKB/WHT). This is
        what lets a non-PANDO uploaded template drive real chart/shape colors
        instead of always falling back to PANDO's own hardcoded palette.

        Nearly every .pptx carries a full 6-accent theme even when nobody ever
        touched it (PowerPoint bakes one in by default), so "the theme has
        >=4 accents" alone is NOT a reliable signal that those accents are the
        deck's real design colors. The reliable signal is whether shapes
        actually REFERENCE the theme via <a:schemeClr val="accentN"/> — that's
        the only way a theme accent shows up in a slide at all, since it's
        never duplicated as literal RGB. If schemeClr accent usage is
        negligible compared to literal RGB fills, the theme is just
        decorative filler and the deck's real palette is the frequent literal
        colors baked directly into its shapes (_dominant_shape_colors) — the
        pattern for hand-designed templates and Keynote/Google Slides exports
        that never wire colors through PowerPoint's theme system at all."""
        accents = [colors.get(f"accent{i}") for i in range(1, 7)]
        accents = [c for c in accents if c]

        scheme_weight, shape_colors = self._scan_color_usage()
        literal_weight = sum(w for _, w in shape_colors)

        use_theme = len(accents) >= 4 and (scheme_weight >= literal_weight or scheme_weight >= 20)
        accent_source = accents if use_theme else [h for h, _ in shape_colors]
        if len(accent_source) < 4:
            return {}

        dark_neutral  = colors.get("dk1") or colors.get("tx1") or "0A231F"
        light_neutral = colors.get("lt2") or colors.get("bg2") or "D9DBD4"
        white         = colors.get("lt1") or colors.get("bg1") or "FFFFFF"
        mapped = {
            "DKG": accent_source[0],
            "MDG": accent_source[1],
            "OLV": accent_source[2],
            "TEL": accent_source[3],
            "LBL": accent_source[4] if len(accent_source) > 4 else light_neutral,
            "GRG": light_neutral,
            "NKB": dark_neutral,
            "WHT": white,
        }
        return {k: v for k, v in mapped.items() if v}

    def _scan_color_usage(self, max_colors: int = 8) -> tuple[int, list[tuple[str, int]]]:
        """Single pass over every real content slide (deliberately NOT layouts
        or masters — those almost always carry decorative theme-accent
        elements baked in by whichever stock PowerPoint design the template
        started from, whether or not the deck's actual slides ever touch
        them, which would otherwise make an unused theme look "used"),
        returning:
          - scheme_weight: how much a shape's own fill/text actually references
            a theme accent color via <a:schemeClr val="accentN"/> (weighted the
            same way as literal fills below) — the only real evidence a theme
            is more than decorative filler, since accent usage never shows up
            as literal RGB.
          - shape_colors: literal RGB fill colors ranked by the same weighting
            (a shape's own fill > a text run's color > a line color), with
            text-highlight annotations dropped entirely (never a brand color),
            near-white/near-black "ink" neutrals dropped (body text is nearly
            always a dark near-black that isn't a real brand accent), and
            near-duplicate hues merged so anti-aliasing doesn't split one
            brand color into several entries."""
        weights: dict[str, int] = {}
        scheme_weight = 0
        try:
            with zipfile.ZipFile(io.BytesIO(self.raw)) as z:
                names = [n for n in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)]
                for name in names:
                    try:
                        root = etree.fromstring(z.read(name))
                    except Exception:
                        continue

                    def _context_weight(node, skip_tags=()):
                        weight, skip = 1, False
                        depth = 0
                        while node is not None and depth < 6:
                            tag = node.tag.split("}")[-1] if isinstance(node.tag, str) and "}" in node.tag else node.tag
                            if tag in skip_tags:
                                skip = True
                                break
                            if tag in ("spPr", "bgPr"):
                                weight = 5
                                break
                            if tag in ("rPr", "defRPr", "endParaRPr"):
                                weight = 2
                                break
                            node = node.getparent()
                            depth += 1
                        return weight, skip

                    for clr in root.iter("{%s}schemeClr" % THEME_NS):
                        if (clr.get("val") or "").startswith("accent"):
                            # p:style's lnRef/fillRef/effectRef are inherited style
                            # hints auto-added to every shape (by PowerPoint and by
                            # python-pptx) that get overridden by the shape's own
                            # explicit spPr fill in virtually all real content —
                            # counting them would make an untouched default theme
                            # look "used" by every single shape.
                            w, skip = _context_weight(clr.getparent(), skip_tags=("highlight", "style"))
                            if not skip:
                                scheme_weight += w

                    for clr in root.iter("{%s}srgbClr" % THEME_NS):
                        val = clr.get("val")
                        if not val or len(val) != 6:
                            continue
                        w, skip = _context_weight(clr.getparent(), skip_tags=("highlight",))
                        if skip:
                            continue
                        val = val.upper()
                        weights[val] = weights.get(val, 0) + w
        except Exception:
            return 0, []

        candidates = [(h, w) for h, w in weights.items() if 40 < _luminance(h) < 245]
        candidates.sort(key=lambda x: -x[1])

        clustered: list[list] = []  # [representative_hex, total_weight]
        for hex_, w in candidates:
            for bucket in clustered:
                if _color_distance(hex_, bucket[0]) < 24:
                    bucket[1] += w
                    break
            else:
                clustered.append([hex_, w])
        clustered.sort(key=lambda b: -b[1])
        return scheme_weight, [(h, w) for h, w in clustered[:max_colors]]

    def _colors(self) -> dict:
        """Extract theme color hex codes from theme1.xml inside the PPTX zip."""
        try:
            with zipfile.ZipFile(io.BytesIO(self.raw)) as z:
                theme_files = [n for n in z.namelist() if "theme/theme" in n]
                if not theme_files:
                    return {}
                xml = z.read(theme_files[0]).decode("utf-8")
            root = ET.fromstring(xml)
            ns = {"a": THEME_NS}
            color_map = {}
            color_scheme = root.find(".//a:clrScheme", ns)
            if color_scheme is None:
                return {}
            for child in color_scheme:
                name = child.tag.split("}")[-1]
                srgb = child.find("a:srgbClr", ns)
                sys_clr = child.find("a:sysClr", ns)
                if srgb is not None:
                    color_map[name] = srgb.get("val", "")
                elif sys_clr is not None:
                    color_map[name] = sys_clr.get("lastClr", "")
            return color_map
        except Exception:
            return {}

    def _fonts(self) -> dict:
        """Extract major/minor font names from theme."""
        try:
            with zipfile.ZipFile(io.BytesIO(self.raw)) as z:
                theme_files = [n for n in z.namelist() if "theme/theme" in n]
                if not theme_files:
                    return {}
                xml = z.read(theme_files[0]).decode("utf-8")
            root = ET.fromstring(xml)
            ns = {"a": THEME_NS}
            fonts = {}
            font_scheme = root.find(".//a:fontScheme", ns)
            if font_scheme is None:
                return {}
            for kind in ["majorFont", "minorFont"]:
                el = font_scheme.find(f"a:{kind}", ns)
                if el is not None:
                    latin = el.find("a:latin", ns)
                    if latin is not None:
                        fonts[kind] = latin.get("typeface", "")
            return fonts
        except Exception:
            return {}

    def _layouts(self) -> list:
        """List available slide masters and layouts with placeholder info."""
        layouts = []
        for mi, master in enumerate(self.prs.slide_masters):
            for li, layout in enumerate(master.slide_layouts):
                phs = {}
                for ph in layout.placeholders:
                    idx = ph.placeholder_format.idx
                    name = PH_NAMES.get(idx, f"ph_{idx}")
                    phs[name] = idx
                layouts.append({
                    "key": f"m{mi}_l{li}",
                    "name": layout.name or f"Layout {li}",
                    "master_idx": mi,
                    "layout_idx": li,
                    "placeholders": phs,
                })
        return layouts
